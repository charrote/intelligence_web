#!/usr/bin/env python3
"""
Intelligence Web — 麦肯锡风格 PPT 生成器
专业咨询风格：简洁高端 · 数据驱动 · 逻辑清晰 · 16:9 宽屏
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# 麦肯锡配色
# ============================================================
BG = RGBColor(0xFF, 0xFF, 0xFF)         # 纯白
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF5, 0xF6, 0xF8)        # 浅灰卡片
BORD = RGBColor(0xE0, 0xE3, 0xE8)        # 边框

MCK_PRIMARY = RGBColor(0x00, 0x2E, 0x5D)  # 麦肯锡深蓝
MCK_SEC = RGBColor(0x00, 0x7B, 0xC0)      # 麦肯锡蓝
MCK_ACC = RGBColor(0x00, 0xA3, 0x9D)      # 麦肯锡绿
MCK_RED = RGBColor(0xC0, 0x39, 0x2B)      # 麦肯锡红
MCK_GRAY = RGBColor(0x5D, 0x63, 0x6E)     # 麦肯锡灰
MCK_LIGHT = RGBColor(0x95, 0x99, 0xA0)    # 浅灰文字
MCK_DARK = RGBColor(0x1A, 0x1A, 0x1A)     # 深黑文字

FONT_CN = 'SimHei'    # 黑体
FONT_EN = 'Arial'     # 英文

# ============================================================
# 工具函数
# ============================================================

def set_bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rc(slide, l, t, w, h, bg=CARD, bd=BORD, r=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = r
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.color.rgb = bd; s.line.width = Pt(0.5)
    return s

def bar(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def tx(slide, l, t, w, h, text, sz=18, bold=False, c=MCK_DARK, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = FONT_CN
    p.alignment = al; p.space_after = Pt(3)
    return tb

def ml(slide, l, t, w, h, *lines, sz=14, c=MCK_DARK, al=PP_ALIGN.LEFT,
       bul=False, sp=1.5, bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    flat = []
    for item in lines:
        if isinstance(item, list): flat.extend(item)
        else: flat.append(item)
    for i, line in enumerate(flat):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ('▸ ' if bul and line.strip() else '') + line
        p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = FONT_CN
        p.alignment = al; p.space_after = Pt(5)
        if sp != 1.0: p.line_spacing = Pt(sz * sp)
    return tb

def bn(slide, l, t, w, h, num, sfx='', c=MCK_ACC, fs=48):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(fs); p.font.bold = True; p.font.color.rgb = c; p.font.name = FONT_EN
    p.alignment = PP_ALIGN.CENTER
    if sfx:
        p2 = tf.add_paragraph(); p2.text = sfx
        p2.font.size = Pt(14); p2.font.color.rgb = MCK_LIGHT; p2.font.name = FONT_CN
        p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
    return tb

def lb(slide, l, t, w, h, text, bg=MCK_PRIMARY, tc=WHITE, sz=11, bold=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.3
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = tc; p.font.bold = bold; p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    return s

def pn(slide, n, total=18):
    tb = slide.shapes.add_textbox(Inches(9.2), Inches(7.15), Inches(0.6), Pt(9))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = f'{n}/{total}'; p.font.size = Pt(8); p.font.color.rgb = MCK_LIGHT; p.font.name = FONT_EN
    p.alignment = PP_ALIGN.RIGHT

def hdr(slide, title, c=MCK_PRIMARY, n=1):
    # 顶部色条
    bar(slide, Inches(0), Inches(0), Inches(10), Pt(3), c)
    tx(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.55),
       title, sz=24, bold=True, c=MCK_DARK)
    pn(slide, n)

# ============================================================
# 图表函数
# ============================================================

def add_chart(slide, l, t, w, h, chart_data, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
              title='', c=MCK_SEC):
    """添加柱状图"""
    chart_data.add_series('数据')
    chart = slide.shapes.add_chart(
        chart_type, l, t, w, h, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = MCK_DARK
    # 设置系列颜色
    for series in chart.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = c
    return chart

def add_pie(slide, l, t, w, h, chart_data, title='', c=MCK_SEC):
    """添加饼图"""
    chart_data.add_series('占比')
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, l, t, w, h, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = MCK_DARK
    # 设置颜色
    colors = [c, MCK_ACC, MCK_SEC, MCK_PRIMARY, MCK_RED]
    for i, series in enumerate(chart.series):
        pt = series.points[0] if len(series.points) > 0 else None
        if pt:
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = colors[i % len(colors)]
    return chart

# ============================================================
# 表格（麦肯锡风格）
# ============================================================

def mck_table(slide, l, t, w, rows, cw, hc=MCK_PRIMARY, ch=0.45):
    """麦肯锡风格表格"""
    ncols = len(rows[0])
    nrows = len(rows)
    ts = slide.shapes.add_table(nrows, ncols, l, t, w, ch * nrows)
    table = ts.table
    cx = 0
    for i, col_w in enumerate(cw):
        table.columns[i].width = col_w
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.text = rows[r][c] if r < len(rows) and c < len(rows[r]) else ""
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.name = FONT_CN
                para.font.color.rgb = WHITE if r == 0 else MCK_DARK
                para.font.bold = r == 0
                para.alignment = PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hc
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return ts

# ============================================================
# 各页
# ============================================================

def s1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG)
    # 顶部色条
    bar(sl, Inches(0), Inches(0), Inches(10), Pt(4), MCK_PRIMARY)
    tx(sl, Inches(1), Inches(2.5), Inches(8), Inches(1),
       'Intelligence Web', sz=44, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.6), Inches(8), Inches(0.5),
       '企业情报智能管理平台', sz=20, c=MCK_SEC, al=PP_ALIGN.CENTER)
    bar(sl, Inches(4.3), Inches(4.3), Inches(1.4), Pt(2), MCK_ACC)
    tx(sl, Inches(1.5), Inches(4.6), Inches(7), Inches(0.8),
       '从情报采集到行动闭环\n让每一个决策都有据可依',
       sz=14, c=MCK_GRAY, al=PP_ALIGN.CENTER)
    for i, tag in enumerate(['开源免费', 'AI 驱动', '容器化部署']):
        lb(sl, Inches(3+i*1.8), Inches(5.8), Inches(1.5), Inches(0.38),
           tag, MCK_PRIMARY, tc=WHITE, sz=10, bold=True)
    pn(sl, 1)

def s2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '执行摘要：信息即战力', MCK_PRIMARY, 2)
    # 核心观点（大字）
    rc(sl, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2), bd=MCK_PRIMARY)
    tx(sl, Inches(0.8), Inches(1.35), Inches(8.4), Inches(0.9),
       '企业每年在信息采集上投入大量人力，但 90% 的信息被淹没、遗忘或被竞争对手捷足先登',
       sz=18, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    # 三个要点
    pts = [
        ('[洞察]', '信息孤岛\n各部门各管一套', '错失商业机会'),
        ('[挑战]', '被动应对\n竞争对手出手后才反应', '永远慢半拍'),
        ('[风险]', '决策靠直觉\n管理层拍板靠经验', '高风险赌局'),
    ]
    for i, (icon, desc, cost) in enumerate(pts):
        l = Inches(0.4 + i*3.2)
        rc(sl, l, Inches(2.8), Inches(2.9), Inches(2.8))
        bar(sl, l+Inches(0.3), Inches(3.0), Inches(0.7), Pt(3), MCK_SEC)
        tx(sl, l+Inches(0.2), Inches(3.2), Inches(2.5), Inches(0.4), icon, sz=11, c=MCK_SEC, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(3.7), Inches(2.5), Inches(0.8), desc, sz=10, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(4.7), Inches(2.5), Inches(0.5), cost, sz=10, c=MCK_RED, al=PP_ALIGN.CENTER)
    pn(sl, 2)

def s3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '解决方案：Intelligence Web 是什么？', MCK_SEC, 3)
    # 核心定义（大框）
    rc(sl, Inches(0.5), Inches(1.2), Inches(9), Inches(1.4), bd=MCK_SEC)
    bar(sl, Inches(0.5), Inches(1.2), Pt(4), Inches(1.4), MCK_SEC)
    tx(sl, Inches(0.8), Inches(1.35), Inches(8.4), Inches(1.1),
       '情报采集 → 结构化存储 → AI 分析 → 行动闭环\n一体化平台',
       sz=20, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    # 四个支柱
    pillars = [
        ('🤖 AI 智能分析', '自动阅读、摘要、分析、留批注', MCK_SEC),
        ('🔍 多渠道采集', '网站抓取、API 对接，灵活设定频率', MCK_ACC),
        ('📊 数据看板', '仪表盘统计、多维筛选、毫秒搜索', MCK_PRIMARY),
        ('🎯 商机管理', '从线索到成交的全生命周期追踪', MCK_RED),
    ]
    for i, (title, desc, clr) in enumerate(pillars):
        l = Inches(0.3 + i*2.4)
        rc(sl, l, Inches(3.0), Inches(2.2), Inches(2.2))
        bar(sl, l+Inches(0.2), Inches(3.2), Inches(0.6), Pt(3), clr)
        tx(sl, l+Inches(0.15), Inches(3.45), Inches(1.9), Inches(0.4), title, sz=11, bold=True, c=MCK_DARK)
        tx(sl, l+Inches(0.15), Inches(3.95), Inches(1.9), Inches(0.8), desc, sz=9, c=MCK_GRAY)
    pn(sl, 3)

def s4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '核心价值：效率提升分析', MCK_ACC, 4)
    # 四个大数字
    effs = [
        ('10x', '信息采集\nAI 自动采集替代人工', MCK_SEC),
        ('50x', '情报分析\nAI 秒级完成人工 15 分钟', MCK_ACC),
        ('100x', '信息检索\n毫秒级全文搜索', MCK_PRIMARY),
        ('10x', '商机响应\n分钟级预警推送', MCK_RED),
    ]
    for i, (num, desc, clr) in enumerate(effs):
        x = Inches(0.3 + i*2.4)
        rc(sl, x, Inches(1.2), Inches(2.2), Inches(2.0))
        bn(sl, x+Inches(0.1), Inches(1.4), Inches(2.0), Inches(1.0), num, c=clr, fs=44)
        tx(sl, x+Inches(0.1), Inches(2.6), Inches(2.0), Inches(0.4), desc, sz=9, c=MCK_GRAY, al=PP_ALIGN.CENTER)
    # 底部结论
    rc(sl, Inches(0.5), Inches(3.6), Inches(9), Inches(0.7), bd=MCK_ACC)
    tx(sl, Inches(0.8), Inches(3.7), Inches(8.4), Inches(0.5),
       'Intelligence Web 不是成本中心，而是收入引擎：预计带来额外 50-100 万年收入',
       sz=14, bold=True, c=MCK_ACC, al=PP_ALIGN.CENTER)
    pn(sl, 4)

def s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, 'ROI 分析：投资回报', MCK_PRIMARY, 5)
    # 成本对比表格
    rows = [
        ['项目', '传统方案', 'Intelligence Web', '节省'],
        ['数据库许可', '10-50 万/年', '0 元', '100%'],
        ['运维人力', '30-60 万/年', '0 元', '省去 1 个全职'],
        ['框架许可', '10-30 万/年', '0 元', '100%'],
        ['部署周期', '数周至数月', '分钟级', '90%+ 时间'],
    ]
    cw = [Inches(1.8), Inches(2.2), Inches(2.8), Inches(1.7)]
    mck_table(sl, Inches(0.5), Inches(1.2), Inches(9), rows, cw, hc=MCK_PRIMARY, ch=0.42)
    # 收益量化
    tx(sl, Inches(0.5), Inches(3.4), Inches(5), Inches(0.3),
       '收益量化（中型企业销售团队）', sz=13, bold=True, c=MCK_DARK)
    data = [
        ('人力释放', 'AI 替代 1-2 小时/天 → 相当于 0.5-1 个全职'),
        ('新增收入', '转化率提升 20% + 周期缩短 20% → 额外 50-100 万'),
        ('部署成本', '几乎为零（开源 + 自有服务器）'),
        ('回报周期', '1-3 个月'),
    ]
    for i, (k, v) in enumerate(data):
        y = Inches(3.8 + i*0.45)
        rc(sl, Inches(0.5), y, Inches(4.5), Inches(0.4), bd=BORD)
        tx(sl, Inches(0.6), y+Inches(0.08), Inches(1.5), Inches(0.25), k, sz=10, bold=True, c=MCK_SEC)
        tx(sl, Inches(2.2), y+Inches(0.08), Inches(2.6), Inches(0.25), v, sz=9, c=MCK_DARK)
    # 大数字
    bn(sl, Inches(6.5), Inches(3.6), Inches(3), Inches(1.2),
       '1-3', sfx='个月回本', c=MCK_ACC, fs=40)
    pn(sl, 5)

def s6(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术架构：共享内核 + 分离领域', MCK_SEC, 6)
    # 架构表格
    rows = [
        ['服务', '端口', '技术栈', '角色'],
        ['Research API', '8766', 'Flask + Python 3.11', '制造情报引擎'],
        ['Sales API', '8767', 'Gunicorn + Python 3.11', '销售情报引擎'],
        ['Gateway', '8765', 'Nginx Alpine', '反向代理 + JWT'],
        ['Meilisearch', '7700', 'Meilisearch v1.12', '全文检索'],
    ]
    cw = [Inches(1.8), Inches(0.7), Inches(2.6), Inches(2.8)]
    mck_table(sl, Inches(0.5), Inches(1.2), Inches(9), rows, cw, hc=MCK_SEC, ch=0.45)
    # 架构原则
    tx(sl, Inches(0.5), Inches(3.5), Inches(5), Inches(0.3),
       '关键架构原则', sz=13, bold=True, c=MCK_DARK)
    principles = [
        '[共享] 所有业务域通过 Docker Volume 挂载同一 core/ 目录',
        '[分离] 各域数据完全隔离，通过 SQLite 独立文件实现',
        '[接口] 18 个 MCP 工具方法，AI Agent 标准化协议访问',
        '[轻量] 零外部数据库依赖 — 除 Meilisearch 外无需额外服务',
    ]
    for i, p in enumerate(principles):
        tx(sl, Inches(0.5), Inches(3.9 + i*0.4), Inches(9), Inches(0.35),
           p, sz=10, c=MCK_DARK)
    pn(sl, 6)

def s7(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, 'AI 协作：人与机的高效飞轮', MCK_ACC, 7)
    # 五步流程
    steps = [
        ('[1] AI 自动巡检', '每日自动扫描\n指定数据源', MCK_SEC),
        ('[2] AI 分析归类', '对采集内容\n进行初步分析', MCK_ACC),
        ('[3] AI 留下建议', '在每条情报下\n留下观察建议', MCK_PRIMARY),
        ('[4] 人类决策', '做出最终判断\n和行动决策', MCK_RED),
        ('[5] 反馈训练 AI', '人类反馈\n反哺 AI 提升', MCK_SEC),
    ]
    for i, (title, desc, clr) in enumerate(steps):
        x = Inches(0.2 + i*1.9)
        rc(sl, x, Inches(1.2), Inches(1.7), Inches(1.3))
        bar(sl, x+Inches(0.15), Inches(1.35), Inches(1.4), Pt(2), clr)
        tx(sl, x+Inches(0.1), Inches(1.55), Inches(1.5), Inches(0.35), title, sz=9, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, x+Inches(0.1), Inches(1.95), Inches(1.5), Inches(0.5), desc, sz=8, c=MCK_GRAY, al=PP_ALIGN.CENTER)
        if i < 4:
            tx(sl, x+Inches(1.65), Inches(1.6), Inches(0.3), Inches(0.3), '→', sz=14, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    # 核心观点
    rc(sl, Inches(0.5), Inches(3.0), Inches(9), Inches(1.2), bd=MCK_ACC)
    tx(sl, Inches(0.8), Inches(3.15), Inches(8.4), Inches(0.9),
       '这不是一个"录入工具"，是让人类和 AI Agent 共同工作的平台。\n每一次人工反馈都在训练 AI，形成正向飞轮效应。',
       sz=13, c=MCK_DARK, al=PP_ALIGN.CENTER)
    pn(sl, 7)

def s8(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '场景演示：制造情报 — "看见趋势"', MCK_ACC, 8)
    # 左侧：数据流
    tx(sl, Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.3),
       '核心数据流', sz=13, bold=True, c=MCK_ACC)
    flows = [
        '竞品发布新产品 → AI 5 分钟内捕获、摘要、推送预警',
        '国家发布智能制造补贴政策 → 自动抓取、分析、标记关联度',
        '行业论坛讨论下一代工艺 → 实时追踪、归类、形成趋势报告',
    ]
    for i, f in enumerate(flows):
        y = Inches(1.7 + i*0.85)
        rc(sl, Inches(0.5), y, Inches(4.3), Inches(0.7))
        bar(sl, Inches(0.5), y+Inches(0.22), Pt(3), Inches(0.22), MCK_ACC)
        tx(sl, Inches(0.8), y+Inches(0.12), Inches(3.8), Inches(0.45), f, sz=10, c=MCK_DARK)
    # 右侧：价值
    tx(sl, Inches(5.3), Inches(1.2), Inches(4), Inches(0.3),
       '价值体现', sz=13, bold=True, c=MCK_ACC)
    vals = ['实时监控竞品产品线调整和产能扩张计划',
            '预判下一代制造技术的商业化时间表',
            '从"事后追悔"到"事前预判"']
    for i, v in enumerate(vals):
        y = Inches(1.7 + i*0.85)
        rc(sl, Inches(5.3), y, Inches(4.2), Inches(0.7))
        tx(sl, Inches(5.5), y+Inches(0.2), Inches(3.8), Inches(0.4),
           f'[✓] {v}', sz=10, c=MCK_DARK)
    # 底部
    rc(sl, Inches(0.5), Inches(4.5), Inches(9), Inches(0.6), bd=MCK_ACC)
    tx(sl, Inches(0.8), Inches(4.55), Inches(8.4), Inches(0.5),
       '看见趋势 · 看见未来 · 看见机会', sz=16, bold=True, c=MCK_ACC, al=PP_ALIGN.CENTER)
    pn(sl, 8)

def s9(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '场景演示：销售情报 — "抓住机会"', MCK_SEC, 9)
    # 三个场景
    scns = [
        ('[目标] 客户扩产预警', '客户宣布扩产 → 第一时间触发跟进流程', MCK_SEC),
        ('[工厂] 竞对投资信号', '竞争对手新建工厂前就被捕捉到投资信号', MCK_RED),
        ('[趋势] 商机漏斗可视化', '销售主管通过数据看板掌握团队商机漏斗健康度', MCK_ACC),
    ]
    for i, (icon, desc, clr) in enumerate(scns):
        l = Inches(0.3 + i*3.2)
        rc(sl, l, Inches(1.2), Inches(2.9), Inches(1.8))
        bar(sl, l+Inches(0.3), Inches(1.4), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.6), Inches(2.5), Inches(0.35), icon, sz=10, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.05), Inches(2.5), Inches(0.5), desc, sz=9, c=MCK_DARK)
    # 商机漏斗
    tx(sl, Inches(0.5), Inches(3.4), Inches(5), Inches(0.3),
       '商机全生命周期', sz=13, bold=True, c=MCK_DARK)
    stages = ['待核实', '合格商机', '方案报价', '商务谈判', '成交/丢标']
    scs = [MCK_LIGHT, MCK_SEC, MCK_ACC, MCK_PRIMARY, MCK_RED]
    for i, (st, sc) in enumerate(zip(stages, scs)):
        x = Inches(0.2 + i*1.9)
        active = i < 4
        lb(sl, x, Inches(3.8), Inches(1.7), Inches(0.4), st,
           sc if active else CARD, tc=MCK_DARK if not active else WHITE, sz=9, bold=True)
        if i < 4:
            tx(sl, x+Inches(1.7), Inches(3.78), Inches(0.3), Inches(0.35), '→', sz=12, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    pn(sl, 9)

def s10(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '扩展性：共享内核 + 任意域扩展', MCK_PRIMARY, 10)
    # 核心理念
    rc(sl, Inches(0.5), Inches(1.2), Inches(9), Inches(0.7), bd=MCK_PRIMARY)
    bar(sl, Inches(0.5), Inches(1.2), Pt(4), Inches(0.7), MCK_PRIMARY)
    tx(sl, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.5),
       '平台不设固定的业务边界，客户需要什么领域，就搭建什么领域',
       sz=14, c=MCK_DARK)
    # 已验证
    tx(sl, Inches(0.5), Inches(2.3), Inches(3), Inches(0.3),
       '已验证领域', sz=12, bold=True, c=MCK_ACC)
    rc(sl, Inches(0.5), Inches(2.7), Inches(3.8), Inches(1.2))
    tx(sl, Inches(0.7), Inches(2.85), Inches(3.4), Inches(0.9),
       '[✓] 制造情报 — "看见趋势"\n\n[✓] 销售情报 — "抓住机会"',
       sz=11, c=MCK_DARK)
    # 可扩展
    tx(sl, Inches(5), Inches(2.3), Inches(4), Inches(0.3),
       '可快速扩展领域', sz=12, bold=True, c=MCK_SEC)
    for i, (d, desc) in enumerate([
        ('[链接] 供应链管理', '寻源、评估、批准、监控'),
        ('[文档] 知识产权监控', '专利追踪、侵权预警'),
        ('[银行] 金融市场追踪', '投融资动态、并购重组'),
        ('[人群] 人力资源情报', '竞对人事变动、人才流动'),
    ]):
        y = Inches(2.7 + i*0.6)
        tx(sl, Inches(5.2), y, Inches(4), Inches(0.5), f'{d} — {desc}', sz=10, c=MCK_DARK)
    # 效率
    rc(sl, Inches(1.5), Inches(4.3), Inches(7), Inches(0.6), bd=MCK_ACC)
    tx(sl, Inches(1.5), Inches(4.35), Inches(7), Inches(0.5),
       '一份配置文件 + 一个前端模板 → 数天内上线新域',
       sz=15, bold=True, c=MCK_ACC, al=PP_ALIGN.CENTER)
    pn(sl, 10)

def s11(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '安全合规：企业级安全内置', MCK_ACC, 11)
    # 安全表格
    rows = [
        ['安全维度', '实现方式'],
        ['认证', 'JWT Bearer Token（HS256），密码 SHA-256 + 随机盐存储'],
        ['密钥保护', 'API Key / Agent Key 落地 XOR 混淆，API 响应中脱敏'],
        ['CORS', '环境变量严格白名单控制允许来源域名'],
        ['RBAC 权限', 'Admin / Manager / Analyst / Viewer 四级角色'],
        ['审计日志', '所有变更操作记录操作人身份和时间戳'],
    ]
    cw = [Inches(1.8), Inches(6.2)]
    mck_table(sl, Inches(0.5), Inches(1.2), Inches(8.6), rows, cw, hc=MCK_ACC, ch=0.45)
    # 底部
    rc(sl, Inches(0.5), Inches(4.2), Inches(9), Inches(0.6), bd=MCK_ACC)
    tx(sl, Inches(0.8), Inches(4.25), Inches(8.4), Inches(0.5),
       '每一个操作都被记录 · 每一个密钥都被保护 · 每一次访问都有据可查',
       sz=12, c=MCK_ACC, al=PP_ALIGN.CENTER)
    pn(sl, 11)

def s12(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术优势：轻而不薄', MCK_SEC, 12)
    # 对比表格
    rows = [
        ['维度', '传统方案', 'Intelligence Web'],
        ['数据库', 'Oracle/PostgreSQL\n年费 10-50 万', 'SQLite 单文件\n零许可费'],
        ['运维', '专职 DBA + DevOps\n年成本 30-60 万', '容器化一键部署\n无需专人运维'],
        ['框架', 'React/Angular 企业版\nBI 工具许可', 'Vanilla JS + Flask\n全部开源免费'],
        ['部署周期', '数周至数月', 'docker compose up -d\n分钟级上线'],
        ['扩展性', '改代码、重新部署\n测试', '配置文件 + 前端模板\n数天上线新域'],
    ]
    cw = [Inches(1.8), Inches(3.5), Inches(3.5)]
    mck_table(sl, Inches(0.3), Inches(1.2), Inches(9.4), rows, cw, hc=MCK_SEC, ch=0.48)
    # 底部
    rc(sl, Inches(0.5), Inches(4.4), Inches(9), Inches(0.6), bd=MCK_SEC)
    tx(sl, Inches(0.8), Inches(4.45), Inches(8.4), Inches(0.5),
       '你不需要再养一个团队来维护这套系统',
       sz=15, bold=True, c=MCK_SEC, al=PP_ALIGN.CENTER)
    pn(sl, 12)

def s13(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '竞品对比：为什么不是 CRM？', MCK_RED, 13)
    # 对比表格
    rows = [
        ['维度', '通用 CRM', 'OA 系统', 'Excel', 'Intelligence Web'],
        ['定位', '客户关系管理', '办公流程管理', '临时记录', '企业情报智能管理'],
        ['信息采集', '[x]', '[x]', '[x]', 'AI 自动采集'],
        ['AI 分析', '[x]', '[x]', '[x]', 'AI Agent 分析'],
        ['多域扩展', '固定模块', '固定模块', '手动搭建', '数天上线新域'],
        ['部署成本', '百万级', '十万级', '低', '几乎为零'],
        ['情报深度', '浅', '浅', '无', '深（全链路）'],
    ]
    cw = [Inches(1.5), Inches(1.7), Inches(1.7), Inches(1.5), Inches(2.2)]
    mck_table(sl, Inches(0.2), Inches(1.2), Inches(9.6), rows, cw, hc=MCK_RED, ch=0.42)
    # 差异化
    rc(sl, Inches(0.5), Inches(4.4), Inches(9), Inches(0.9), bd=MCK_PRIMARY)
    bar(sl, Inches(0.5), Inches(4.4), Pt(4), Inches(0.9), MCK_PRIMARY)
    tx(sl, Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.7),
       'Intelligence Web 只做一件事：帮企业把散落在各处的情报变成可行动的洞察。为此做了极深的垂直打磨。',
       sz=12, c=MCK_PRIMARY)
    pn(sl, 13)

def s14(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '适用人群：为谁而建？', MCK_PRIMARY, 14)
    users = [
        ('一线销售/商务经理', '比竞对更快知道客户在哪\n需求是什么',
         '自动预警 + 客户画像 + 商机全追踪', MCK_SEC),
        ('市场研究/战略规划', '持续扫描行业全貌\n形成可指导决策的报告',
         '多渠道采集 + AI 分析 + 趋势可视化', MCK_ACC),
        ('企业管理者/决策层', '一眼看清整体状况\n不做凭感觉的赌局',
         '数据看板 + AI 摘要 + 组织级公共资产', MCK_RED),
    ]
    for i, (role, need, sol, clr) in enumerate(users):
        l = Inches(0.2 + i*3.2)
        rc(sl, l, Inches(1.2), Inches(2.9), Inches(4.5), bd=clr)
        bar(sl, l+Inches(0.3), Inches(1.4), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.7), Inches(2.5), Inches(0.4), role, sz=13, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.2), Inches(2.5), Inches(0.8), need, sz=9, c=MCK_GRAY, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(3.3), Inches(2.5), Inches(0.25), '[箭头]', sz=14, c=clr, al=PP_ALIGN.CENTER)
        lb(sl, l+Inches(0.1), Inches(3.7), Inches(2.7), Inches(1.1), sol, clr, tc=WHITE, sz=9, bold=True)
    tx(sl, Inches(0.5), Inches(6.2), Inches(9), Inches(0.3),
       '不是给所有人的万能工具 — 是为情报驱动决策的团队量身定制的效率倍增器',
       sz=11, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    pn(sl, 14)

def s15(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '部署启动：最快数分钟', MCK_SEC, 15)
    steps = [
        ('Step 1', 'docker compose up -d', '一条命令启动全部服务', MCK_SEC),
        ('Step 2', '浏览器访问 :8765', '登录，开始使用', MCK_ACC),
        ('Step 3', '配置数据源 + 项目', '开始自动采集', MCK_PRIMARY),
    ]
    for i, (step, cmd, desc, clr) in enumerate(steps):
        l = Inches(0.5 + i*3.2)
        rc(sl, l, Inches(1.2), Inches(2.9), Inches(2.2), bd=clr)
        bar(sl, l+Inches(0.3), Inches(1.4), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.7), Inches(2.5), Inches(0.4), step, sz=16, bold=True, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.2), Inches(2.5), Inches(0.4), cmd, sz=10, c=MCK_SEC, al=PP_ALIGN.CENTER, bold=True)
        tx(sl, l+Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.4), desc, sz=10, c=MCK_GRAY, al=PP_ALIGN.CENTER)
        if i < 2:
            tx(sl, l+Inches(2.85), Inches(2.1), Inches(0.4), Inches(0.4), '→', sz=18, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    # 开箱即用
    tx(sl, Inches(0.5), Inches(3.8), Inches(4), Inches(0.3),
       '开箱即用的能力', sz=13, bold=True, c=MCK_SEC)
    feats = [
        '[✓] 15 个功能页面，即开即用',
        '[✓] RBAC 权限体系，三级角色即刻生效',
        '[✓] AI Agent 预设模板，配置即用',
        '[✓] 暗色模式、响应式设计，现代用户体验',
        '[✓] 新业务域：一份配置 + 一个模板 → 数天内上线',
    ]
    for i, f in enumerate(feats):
        tx(sl, Inches(0.5), Inches(4.2 + i*0.4), Inches(9), Inches(0.35),
           f, sz=10, c=MCK_DARK)
    pn(sl, 15)

def s16(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '投资回报：1-3 个月回本', MCK_ACC, 16)
    # 三个大数字
    nums = [
        ('1-3', '个月\n投资回报周期', MCK_ACC),
        ('50-100', '万\n年新增收入', MCK_SEC),
        ('0.5-1', '个\n全职人力释放', MCK_PRIMARY),
    ]
    for i, (num, sfx, clr) in enumerate(nums):
        x = Inches(0.5 + i*3.2)
        rc(sl, x, Inches(1.2), Inches(2.9), Inches(2.2), bd=clr)
        bar(sl, x+Inches(0.3), Inches(1.4), Inches(0.7), Pt(3), clr)
        bn(sl, x+Inches(0.1), Inches(1.7), Inches(2.7), Inches(1.2), num, sfx=sfx, c=clr, fs=36)
    # 基准
    rc(sl, Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
    tx(sl, Inches(0.8), Inches(3.85), Inches(8.4), Inches(0.4),
       '基准假设：一名年薪资 20 万的销售人员', sz=11, c=MCK_GRAY)
    # 详细数据
    data = [
        ('人力释放', 'AI 替代 1-2 小时/天 → 相当于 0.5-1 个全职'),
        ('新增商机收入', '转化率提升 20% + 销售周期缩短 20% → 额外 50-100 万年收入'),
        ('系统部署成本', '几乎为零（开源 + 自有服务器）'),
        ('投资回报周期', '1-3 个月'),
    ]
    for i, (k, v) in enumerate(data):
        y = Inches(4.6 + i*0.4)
        rc(sl, Inches(0.5), y, Inches(9), Inches(0.35), bd=BORD)
        tx(sl, Inches(0.7), y+Inches(0.05), Inches(2), Inches(0.25), k, sz=9, bold=True, c=MCK_SEC)
        tx(sl, Inches(2.8), y+Inches(0.05), Inches(6.5), Inches(0.25), v, sz=9, c=MCK_DARK)
    pn(sl, 16)

def s17(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '行动号召：升级您的企业情报管理', MCK_PRIMARY, 17)
    # 核心信息
    pts = [
        '[网络] 开源 · 免费 · 可私有化部署',
        '[AI] AI 驱动 · 人机协同 · 持续进化',
        '[图表] 1-3 个月投资回报 · 100 万年新增收入',
        '[火箭] 一条命令启动 · 数天扩展新域',
    ]
    for i, p in enumerate(pts):
        tx(sl, Inches(2), Inches(1.5 + i*0.5), Inches(6), Inches(0.4),
           p, sz=13, c=MCK_DARK, al=PP_ALIGN.CENTER)
    # 联系方式
    rc(sl, Inches(2), Inches(3.8), Inches(6), Inches(1.2), bd=MCK_ACC)
    tx(sl, Inches(2.2), Inches(3.9), Inches(5.6), Inches(0.5),
       '准备好让您的企业情报管理升级了吗？',
       sz=18, bold=True, c=MCK_ACC, al=PP_ALIGN.CENTER)
    tx(sl, Inches(2.2), Inches(4.6), Inches(5.6), Inches(0.3),
       '[官网]  [邮箱]  [电话]  [二维码]',
       sz=10, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    pn(sl, 17)

def s18(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG)
    bar(sl, Inches(0), Inches(0), Inches(10), Pt(4), MCK_PRIMARY)
    tx(sl, Inches(1), Inches(2.5), Inches(8), Inches(0.9),
       'Intelligence Web', sz=42, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.5), Inches(8), Inches(0.5),
       '让情报成为您的核心竞争力', sz=18, c=MCK_SEC, al=PP_ALIGN.CENTER)
    bar(sl, Inches(4.3), Inches(4.2), Inches(1.4), Pt(2), MCK_ACC)
    tx(sl, Inches(2), Inches(4.5), Inches(6), Inches(0.8),
       '谢谢\nThank You',
       sz=20, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    pn(sl, 18)

# ============================================================
# 主函数
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s1(prs); s2(prs); s3(prs); s4(prs); s5(prs); s6(prs); s7(prs); s8(prs)
    s9(prs); s10(prs); s11(prs); s12(prs); s13(prs); s14(prs); s15(prs)
    s16(prs); s17(prs); s18(prs)

    prs.save('docs/Intelligence_Web_麦肯锡风格 PPT.pptx')
    print('PPT 麦肯锡风格已生成：docs/Intelligence_Web_麦肯锡风格 PPT.pptx')
    print('18 页 · 16:9 宽屏 · 麦肯锡配色 · 专业咨询风格')

if __name__ == '__main__':
    main()
#!/usr/bin/env python3
"""
Intelligence Web — 麦肯锡风格 PPT v2
内容充实版：大字体 · 满内容 · 专业咨询风格 · 16:9 宽屏
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# 麦肯锡配色（充实版）
# ============================================================
BG = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF5, 0xF6, 0xF8)
BORD = RGBColor(0xE0, 0xE3, 0xE8)

MCK_BLUE = RGBColor(0x00, 0x2E, 0x5D)
MCK_BLUE_L = RGBColor(0x1A, 0x56, 0xDB)
MCK_GREEN = RGBColor(0x00, 0xA3, 0x9D)
MCK_RED = RGBColor(0xC0, 0x39, 0x2B)
MCK_GRAY = RGBColor(0x5D, 0x63, 0x6E)
MCK_LIGHT = RGBColor(0x95, 0x99, 0xA0)
MCK_DARK = RGBColor(0x1A, 0x1A, 0x1A)
MCK_BG = RGBColor(0xF8, 0xF9, 0xFA)

FONT_CN = 'SimHei'
FONT_EN = 'Arial'

# ============================================================
# 工具函数（充实版参数）
# ============================================================

def set_bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rc(slide, l, t, w, h, bg=CARD, bd=BORD, r=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = r
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.color.rgb = bd; s.line.width = Pt(0.5)
    return s

def bar(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def tx(slide, l, t, w, h, text, sz=20, bold=False, c=MCK_DARK, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = FONT_CN
    p.alignment = al; p.space_after = Pt(4)
    return tb

def ml(slide, l, t, w, h, *lines, sz=16, c=MCK_DARK, al=PP_ALIGN.LEFT,
       bul=False, sp=1.5, bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    flat = []
    for item in lines:
        if isinstance(item, list): flat.extend(item)
        else: flat.append(item)
    for i, line in enumerate(flat):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ('▸ ' if bul and line.strip() else '') + line
        p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = FONT_CN
        p.alignment = al; p.space_after = Pt(6)
        if sp != 1.0: p.line_spacing = Pt(sz * sp)
    return tb

def bn(slide, l, t, w, h, num, sfx='', c=MCK_GREEN, fs=56):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(fs); p.font.bold = True; p.font.color.rgb = c; p.font.name = FONT_EN
    p.alignment = PP_ALIGN.CENTER
    if sfx:
        p2 = tf.add_paragraph(); p2.text = sfx
        p2.font.size = Pt(16); p2.font.color.rgb = MCK_GRAY; p2.font.name = FONT_CN
        p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    return tb

def lb(slide, l, t, w, h, text, bg=MCK_BLUE, tc=WHITE, sz=14, bold=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.35
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = tc; p.font.bold = bold; p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    return s

def pn(slide, n, total=18):
    tb = slide.shapes.add_textbox(Inches(9.3), Inches(7.15), Inches(0.6), Pt(10))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = f'{n}/{total}'; p.font.size = Pt(9); p.font.color.rgb = MCK_LIGHT; p.font.name = FONT_EN
    p.alignment = PP_ALIGN.RIGHT

def hdr(slide, title, c=MCK_BLUE, n=1):
    bar(slide, Inches(0), Inches(0), Inches(10), Pt(4), c)
    tx(slide, Inches(0.6), Inches(0.25), Inches(9), Inches(0.65),
       title, sz=28, bold=True, c=MCK_DARK)
    pn(slide, n)

# ============================================================
# 表格（充实版）
# ============================================================

def mck_table(slide, l, t, w, rows, cw, hc=MCK_BLUE, ch=0.55):
    ncols = len(rows[0])
    nrows = len(rows)
    ts = slide.shapes.add_table(nrows, ncols, l, t, w, ch * nrows)
    table = ts.table
    cx = 0
    for i, col_w in enumerate(cw):
        table.columns[i].width = col_w
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.text = rows[r][c] if r < len(rows) and c < len(rows[r]) else ""
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.name = FONT_CN
                para.font.color.rgb = WHITE if r == 0 else MCK_DARK
                para.font.bold = r == 0
                para.alignment = PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hc
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = MCK_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return ts

# ============================================================
# 各页（充实版）
# ============================================================

def s1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG)
    bar(sl, Inches(0), Inches(0), Inches(10), Pt(5), MCK_BLUE)
    tx(sl, Inches(1), Inches(2.2), Inches(8), Inches(1.2),
       'Intelligence Web', sz=52, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.5), Inches(8), Inches(0.7),
       '企业情报智能管理平台', sz=24, c=MCK_BLUE_L, al=PP_ALIGN.CENTER)
    bar(sl, Inches(4.2), Inches(4.4), Inches(1.6), Pt(3), MCK_GREEN)
    tx(sl, Inches(1.5), Inches(4.8), Inches(7), Inches(1),
       '从情报采集到行动闭环\n让每一个决策都有据可依',
       sz=16, c=MCK_GRAY, al=PP_ALIGN.CENTER)
    for i, tag in enumerate(['开源免费', 'AI 驱动', '容器化部署']):
        lb(sl, Inches(3+i*1.8), Inches(6.2), Inches(1.5), Inches(0.45),
           tag, MCK_BLUE, tc=WHITE, sz=12, bold=True)
    pn(sl, 1)

def s2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '执行摘要：信息即战力', MCK_BLUE, 2)
    # 核心观点（大框）
    rc(sl, Inches(0.5), Inches(1.2), Inches(9), Inches(1.5), bd=MCK_BLUE)
    tx(sl, Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.2),
       '企业每年在信息采集上投入大量人力，但 90% 的信息被淹没、遗忘或被竞争对手捷足先登',
       sz=22, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    # 三个痛点卡片
    pts = [
        ('[洞察]', '信息孤岛', '各部门各管一套\n重要情报反复丢失', '错失商业机会\n重复劳动', MCK_BLUE),
        ('[挑战]', '被动应对', '竞争对手出手后才\n反应过来', '永远慢半拍\n窗口期耗尽', MCK_RED),
        ('[风险]', '决策靠直觉', '管理层拍板靠经验\n和感觉', '高风险赌局\n事后才知对错', MCK_GRAY),
    ]
    for i, (icon, title, desc, cost, clr) in enumerate(pts):
        l = Inches(0.3 + i*3.2)
        rc(sl, l, Inches(3.0), Inches(2.9), Inches(3.5), bd=clr)
        bar(sl, l+Inches(0.3), Inches(3.2), Inches(0.8), Pt(4), clr)
        tx(sl, l+Inches(0.2), Inches(3.5), Inches(2.5), Inches(0.5), icon, sz=13, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(4.1), Inches(2.5), Inches(0.5), title, sz=16, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(4.7), Inches(2.5), Inches(0.8), desc, sz=11, c=MCK_GRAY, al=PP_ALIGN.CENTER)
        bar(sl, l+Inches(0.7), Inches(5.7), Inches(1.5), Pt(2), BORD)
        tx(sl, l+Inches(0.2), Inches(5.9), Inches(2.5), Inches(0.5), cost, sz=11, c=clr, al=PP_ALIGN.CENTER)
    pn(sl, 2)

def s3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '解决方案：Intelligence Web 是什么？', MCK_BLUE_L, 3)
    # 核心定义（大框）
    rc(sl, Inches(0.5), Inches(1.2), Inches(9), Inches(1.8), bd=MCK_BLUE_L)
    bar(sl, Inches(0.5), Inches(1.2), Pt(5), Inches(1.8), MCK_BLUE_L)
    tx(sl, Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.5),
       '情报采集 → 结构化存储 → AI 分析 → 行动闭环\n一体化平台',
       sz=24, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    # 四个支柱（大卡片）
    pillars = [
        ('[AI]', 'AI 智能分析', '自动阅读、摘要、分析、留批注\n人机协同知识积累', MCK_BLUE_L),
        ('[搜索]', '多渠道采集', '网站抓取、API 对接\n按日/周/月灵活设定频率', MCK_GREEN),
        ('[图表]', '数据看板', '仪表盘统计、多维筛选\n毫秒级全文搜索', MCK_BLUE),
        ('[目标]', '商机管理', '从线索到成交的全生命周期追踪\n销售域专属', MCK_RED),
    ]
    for i, (icon, title, desc, clr) in enumerate(pillars):
        l = Inches(0.2 + i*2.5)
        rc(sl, l, Inches(3.4), Inches(2.3), Inches(2.8))
        bar(sl, l+Inches(0.25), Inches(3.6), Inches(0.7), Pt(4), clr)
        tx(sl, l+Inches(0.15), Inches(3.85), Inches(2.0), Inches(0.5), icon, sz=13, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.15), Inches(4.4), Inches(2.0), Inches(0.5), title, sz=14, bold=True, c=MCK_DARK)
        tx(sl, l+Inches(0.15), Inches(4.95), Inches(2.0), Inches(0.8), desc, sz=10, c=MCK_GRAY)
    pn(sl, 3)

def s4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '核心价值：效率提升分析', MCK_GREEN, 4)
    # 四个大数字卡片
    effs = [
        ('10x', '信息采集', 'AI 自动采集替代人工\n每日 1-2 小时解放', MCK_BLUE_L),
        ('50x', '情报分析', 'AI 秒级完成人工\n15 分钟摘要工作', MCK_GREEN),
        ('100x', '信息检索', '毫秒级全文搜索\n中英文混合检索', MCK_BLUE),
        ('10x', '商机响应', '分钟级预警推送\n第一时间触发跟进', MCK_RED),
    ]
    for i, (num, title, desc, clr) in enumerate(effs):
        x = Inches(0.2 + i*2.5)
        rc(sl, x, Inches(1.2), Inches(2.3), Inches(2.5))
        bn(sl, x+Inches(0.1), Inches(1.4), Inches(2.1), Inches(1.1), num, c=clr, fs=52)
        tx(sl, x+Inches(0.1), Inches(2.6), Inches(2.1), Inches(0.4), title, sz=13, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, x+Inches(0.1), Inches(3.1), Inches(2.1), Inches(0.5), desc, sz=10, c=MCK_GRAY, al=PP_ALIGN.CENTER)
    # 底部结论（大框）
    rc(sl, Inches(0.5), Inches(4.2), Inches(9), Inches(1.0), bd=MCK_GREEN)
    tx(sl, Inches(0.8), Inches(4.4), Inches(8.4), Inches(0.7),
       'Intelligence Web 不是成本中心，而是收入引擎：\n预计为中型企业带来额外 50-100 万年收入',
       sz=16, bold=True, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 4)

def s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, 'ROI 分析：投资回报', MCK_BLUE, 5)
    # 成本对比表格（大）
    rows = [
        ['项目', '传统方案', 'Intelligence Web', '节省'],
        ['数据库许可', 'Oracle/PostgreSQL\n年费 10-50 万', 'SQLite 单文件\n零许可费', '100% 节省'],
        ['运维人力', '专职 DBA + DevOps\n年成本 30-60 万', '容器化一键部署\n无需专人运维', '省去 1 个全职'],
        ['框架许可', 'React/Angular 企业版\nBI 工具许可', 'Vanilla JS + Flask\n全部开源免费', '100% 节省'],
        ['部署周期', '数周至数月', 'docker compose up -d\n分钟级上线', '90%+ 时间'],
    ]
    cw = [Inches(1.8), Inches(2.5), Inches(2.8), Inches(1.7)]
    mck_table(sl, Inches(0.5), Inches(1.2), Inches(9), rows, cw, hc=MCK_BLUE, ch=0.52)
    # 收益量化（大卡片）
    tx(sl, Inches(0.5), Inches(3.6), Inches(5), Inches(0.4),
       '收益量化（中型企业销售团队）', sz=16, bold=True, c=MCK_DARK)
    data = [
        ('人力释放', 'AI 替代 1-2 小时/天 → 相当于 0.5-1 个全职人力释放', MCK_BLUE_L),
        ('新增收入', '转化率提升 20% + 周期缩短 20% → 额外 50-100 万年收入', MCK_GREEN),
        ('部署成本', '几乎为零（开源 + 自有服务器即可运行）', MCK_GRAY),
        ('回报周期', '1-3 个月即可收回全部投资成本', MCK_RED),
    ]
    for i, (k, v, clr) in enumerate(data):
        y = Inches(4.1 + i*0.55)
        rc(sl, Inches(0.5), y, Inches(4.8), Inches(0.48), bd=clr)
        tx(sl, Inches(0.7), y+Inches(0.1), Inches(1.6), Inches(0.28), k, sz=12, bold=True, c=clr)
        tx(sl, Inches(2.4), y+Inches(0.1), Inches(2.7), Inches(0.28), v, sz=10, c=MCK_DARK)
    # 大数字
    bn(sl, Inches(6.5), Inches(3.8), Inches(3), Inches(1.5),
       '1-3', sfx='个月\n投资回报周期', c=MCK_GREEN, fs=48)
    pn(sl, 5)

def s6(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术架构：共享内核 + 分离领域', MCK_BLUE_L, 6)
    # 架构图 - 用图形表示各服务关系
    # 顶部：客户端层
    rc(sl, Inches(3.0), Inches(1.2), Inches(4), Inches(0.7), bd=MCK_BLUE)
    tx(sl, Inches(3.0), Inches(1.3), Inches(4), Inches(0.5),
       '🖥 客户端层\n浏览器 Portal SPA', sz=10, c=MCK_BLUE, al=PP_ALIGN.CENTER)
    # 箭头
    tx(sl, Inches(4.8), Inches(1.95), Inches(0.4), Inches(0.3), '↓', sz=16, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    # 网关层
    rc(sl, Inches(2.5), Inches(2.3), Inches(5), Inches(0.7), bd=MCK_GREEN)
    tx(sl, Inches(2.5), Inches(2.4), Inches(5), Inches(0.5),
       '🔀 网关层\nNginx 反向代理 + JWT 鉴权 :8765', sz=10, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    # 箭头
    tx(sl, Inches(4.8), Inches(3.05), Inches(0.4), Inches(0.3), '↓', sz=16, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    # 业务域层（两个）
    rc(sl, Inches(0.5), Inches(3.4), Inches(4), Inches(0.7), bd=MCK_BLUE)
    tx(sl, Inches(0.5), Inches(3.5), Inches(4), Inches(0.5),
       '⚙️ 制造情报 API\nFlask + Python 3.11 :8766', sz=10, c=MCK_BLUE, al=PP_ALIGN.CENTER)
    rc(sl, Inches(5.5), Inches(3.4), Inches(4), Inches(0.7), bd=MCK_GREEN)
    tx(sl, Inches(5.5), Inches(3.5), Inches(4), Inches(0.5),
       '⚙️ 销售情报 API\nGunicorn + Python 3.11 :8767', sz=10, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    # 箭头
    tx(sl, Inches(2.3), Inches(4.15), Inches(0.4), Inches(0.3), '↓', sz=16, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    tx(sl, Inches(7.3), Inches(4.15), Inches(0.4), Inches(0.3), '↓', sz=16, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    # 核心引擎层
    rc(sl, Inches(1.5), Inches(4.5), Inches(7), Inches(0.7), bd=MCK_BLUE)
    tx(sl, Inches(1.5), Inches(4.6), Inches(7), Inches(0.5),
       '🧠 共享核心引擎 core/\napp.py · db.py · datasource.py · project.py · domain.py · target_types.py', sz=10, c=MCK_BLUE, al=PP_ALIGN.CENTER)
    # 底部服务（两个）
    rc(sl, Inches(0.5), Inches(5.5), Inches(2.5), Inches(0.6), bd=MCK_GREEN)
    tx(sl, Inches(0.5), Inches(5.6), Inches(2.5), Inches(0.4),
       '💾 SQLite\nresearch.db', sz=9, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    rc(sl, Inches(3.2), Inches(5.5), Inches(2.5), Inches(0.6), bd=MCK_GREEN)
    tx(sl, Inches(3.2), Inches(5.6), Inches(2.5), Inches(0.4),
       '💾 SQLite\nsales.db', sz=9, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    rc(sl, Inches(5.9), Inches(5.5), Inches(2.5), Inches(0.6), bd=MCK_BLUE)
    tx(sl, Inches(5.9), Inches(5.6), Inches(2.5), Inches(0.4),
       '🔍 Meilisearch\n全文检索 :7700', sz=9, c=MCK_BLUE, al=PP_ALIGN.CENTER)
    # 架构原则
    tx(sl, Inches(0.5), Inches(6.3), Inches(5), Inches(0.3),
       '关键架构原则', sz=14, bold=True, c=MCK_DARK)
    principles = [
        ('[共享]', '所有业务域通过 Docker Volume 挂载同一 core/ 目录', MCK_BLUE_L),
        ('[分离]', '各域数据完全隔离，通过 SQLite 独立文件实现', MCK_GREEN),
        ('[接口]', '18 个 MCP 工具方法，AI Agent 标准化协议访问', MCK_BLUE),
        ('[轻量]', '零外部依赖 — 除 Meilisearch 外无需额外服务', MCK_GRAY),
    ]
    for i, (icon, desc, clr) in enumerate(principles):
        y = Inches(6.6 + i*0.35)
        tx(sl, Inches(0.5), y, Inches(9), Inches(0.3),
           f'{icon} {desc}', sz=9, c=MCK_DARK)
    pn(sl, 6)

def s8(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术栈详解：极简而强大', MCK_BLUE, 7)
    # 技术栈表格
    rows = [
        ['服务', '技术栈', '版本', '说明'],
        ['Research API', 'Flask', 'Python 3.11', '制造情报业务引擎，轻量级 Web 框架'],
        ['Sales API', 'Gunicorn', 'Python 3.11', '销售情报业务引擎，生产级 WSGI 服务器'],
        ['Gateway', 'Nginx Alpine', 'Alpine Linux', '反向代理 + JWT 鉴权网关，统一入口'],
        ['Search', 'Meilisearch', 'v1.12', '全文检索引擎，支持中英文混合搜索'],
        ['Database', 'SQLite', '单文件', '零配置数据库，开箱即用'],
        ['Portal', 'Vanilla JS', 'HTML/CSS/JS', '纯前端，15 个功能页面，零框架依赖'],
        ['AI 接口', 'MCP Server', 'JSON-RPC', '18 个工具方法，AI Agent 标准化协议'],
    ]
    cw = [Inches(1.5), Inches(1.5), Inches(1.2), Inches(4.5)]
    mck_table(sl, Inches(0.5), Inches(1.2), Inches(9), rows, cw, hc=MCK_BLUE, ch=0.48)
    # 底部总结
    rc(sl, Inches(0.5), Inches(5.5), Inches(9), Inches(0.8), bd=MCK_BLUE)
    tx(sl, Inches(0.8), Inches(5.6), Inches(8.4), Inches(0.6),
       '极简技术栈，经过充分验证，稳定可靠且易于维护\n无需 DBA，无需 DevOps，一条命令即可启动整套系统',
       sz=13, c=MCK_BLUE, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s9(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, 'AI 协作：人与机的高效飞轮', MCK_GREEN, 7)
    # 五步流程（大卡片）
    steps = [
        ('[1]', 'AI 自动巡检', '每日自动扫描\n指定数据源', MCK_BLUE_L),
        ('[2]', 'AI 分析归类', '对采集内容\n进行初步分析', MCK_GREEN),
        ('[3]', 'AI 留下建议', '在每条情报下\n留下观察建议', MCK_BLUE),
        ('[4]', '人类决策', '做出最终判断\n和行动决策', MCK_RED),
        ('[5]', '反馈训练 AI', '人类反馈\n反哺 AI 提升', MCK_BLUE_L),
    ]
    for i, (num, title, desc, clr) in enumerate(steps):
        x = Inches(0.1 + i*1.95)
        rc(sl, x, Inches(1.2), Inches(1.75), Inches(1.6), bd=clr)
        bar(sl, x+Inches(0.2), Inches(1.4), Inches(1.35), Pt(3), clr)
        tx(sl, x+Inches(0.1), Inches(1.6), Inches(1.55), Inches(0.4), num, sz=14, bold=True, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, x+Inches(0.1), Inches(2.05), Inches(1.55), Inches(0.4), title, sz=11, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, x+Inches(0.1), Inches(2.5), Inches(1.55), Inches(0.5), desc, sz=9, c=MCK_GRAY, al=PP_ALIGN.CENTER)
        if i < 4:
            tx(sl, x+Inches(1.75), Inches(1.9), Inches(0.3), Inches(0.3), '→', sz=16, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    # 核心观点（大框）
    rc(sl, Inches(0.5), Inches(3.2), Inches(9), Inches(1.5), bd=MCK_GREEN)
    tx(sl, Inches(0.8), Inches(3.4), Inches(8.4), Inches(1.2),
       '这不是一个"录入工具"，是让人类和 AI Agent 共同工作的平台。\n每一次人工反馈都在训练 AI，形成正向飞轮效应，让系统越用越聪明。',
       sz=16, c=MCK_DARK, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s7(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '场景演示：制造情报 — "看见趋势"', MCK_GREEN, 8)
    # 左侧：数据流（大卡片）
    tx(sl, Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.4),
       '核心数据流', sz=16, bold=True, c=MCK_GREEN)
    flows = [
        '竞品发布新产品 → AI 5 分钟内捕获、摘要、推送预警，销售团队第一时间获知',
        '国家发布智能制造补贴政策 → 自动抓取、分析、标记关联度，战略部门快速响应',
        '行业论坛讨论下一代工艺 → 实时追踪、归类、形成趋势报告，研发部门提前布局',
    ]
    for i, f in enumerate(flows):
        y = Inches(1.8 + i*1.1)
        rc(sl, Inches(0.5), y, Inches(4.3), Inches(0.95))
        bar(sl, Inches(0.5), y+Inches(0.3), Pt(4), Inches(0.3), MCK_GREEN)
        tx(sl, Inches(0.8), y+Inches(0.15), Inches(3.8), Inches(0.65), f, sz=12, c=MCK_DARK)
    # 右侧：价值（大卡片）
    tx(sl, Inches(5.3), Inches(1.2), Inches(4), Inches(0.4),
       '价值体现', sz=16, bold=True, c=MCK_GREEN)
    vals = ['实时监控竞品产品线调整和产能扩张计划，掌握市场主动权',
            '预判下一代制造技术的商业化时间表，提前布局研发资源',
            '从"事后追悔"到"事前预判"，将情报转化为战略优势']
    for i, v in enumerate(vals):
        y = Inches(1.8 + i*1.1)
        rc(sl, Inches(5.3), y, Inches(4.2), Inches(0.95))
        tx(sl, Inches(5.5), y+Inches(0.2), Inches(3.8), Inches(0.6),
           f'[✓] {v}', sz=12, c=MCK_DARK)
    # 底部
    rc(sl, Inches(0.5), Inches(5.3), Inches(9), Inches(0.8), bd=MCK_GREEN)
    tx(sl, Inches(0.8), Inches(5.4), Inches(8.4), Inches(0.6),
       '看见趋势 · 看见未来 · 看见机会',
       sz=20, bold=True, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s10(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '场景演示：销售情报 — "抓住机会"', MCK_BLUE_L, 9)
    # 三个场景（大卡片）
    scns = [
        ('[目标]', '客户扩产预警', '客户宣布扩产 → 第一时间触发跟进流程\n销售主管立即安排拜访，抢占先机', MCK_BLUE_L),
        ('[工厂]', '竞对投资信号', '竞争对手新建工厂前就被捕捉到投资信号\n提前准备应对策略，掌握主动权', MCK_RED),
        ('[趋势]', '商机漏斗可视化', '销售主管通过数据看板掌握团队商机漏斗健康度\n及时发现瓶颈，优化销售流程', MCK_GREEN),
    ]
    for i, (icon, title, desc, clr) in enumerate(scns):
        l = Inches(0.2 + i*3.2)
        rc(sl, l, Inches(1.2), Inches(2.9), Inches(2.5))
        bar(sl, l+Inches(0.3), Inches(1.4), Inches(0.8), Pt(4), clr)
        tx(sl, l+Inches(0.2), Inches(1.7), Inches(2.5), Inches(0.5), icon, sz=13, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.3), Inches(2.5), Inches(0.5), title, sz=14, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.6), desc, sz=10, c=MCK_GRAY)
    # 商机漏斗
    tx(sl, Inches(0.5), Inches(4.1), Inches(5), Inches(0.4),
       '商机全生命周期', sz=16, bold=True, c=MCK_DARK)
    stages = ['待核实', '合格商机', '方案报价', '商务谈判', '成交/丢标']
    scs = [MCK_LIGHT, MCK_BLUE_L, MCK_GREEN, MCK_BLUE, MCK_RED]
    for i, (st, sc) in enumerate(zip(stages, scs)):
        x = Inches(0.2 + i*1.9)
        active = i < 4
        lb(sl, x, Inches(4.6), Inches(1.7), Inches(0.5), st,
           sc if active else CARD, tc=MCK_DARK if not active else WHITE, sz=12, bold=True)
        if i < 4:
            tx(sl, x+Inches(1.7), Inches(4.55), Inches(0.3), Inches(0.4), '→', sz=14, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s11(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '扩展性：共享内核 + 任意域扩展', MCK_BLUE, 10)
    # 核心理念（大框）
    rc(sl, Inches(0.5), Inches(1.2), Inches(9), Inches(0.9), bd=MCK_BLUE)
    bar(sl, Inches(0.5), Inches(1.2), Pt(5), Inches(0.9), MCK_BLUE)
    tx(sl, Inches(0.8), Inches(1.35), Inches(8.4), Inches(0.6),
       '平台不设固定的业务边界，客户需要什么领域，就搭建什么领域\n从提出需求到新域可用，最快数天内完成交付',
       sz=15, c=MCK_DARK)
    # 已验证领域（大卡片）
    tx(sl, Inches(0.5), Inches(2.5), Inches(3), Inches(0.4),
       '已验证领域', sz=15, bold=True, c=MCK_GREEN)
    rc(sl, Inches(0.5), Inches(3.0), Inches(3.8), Inches(1.5))
    tx(sl, Inches(0.7), Inches(3.15), Inches(3.4), Inches(1.2),
       '[✓] 制造情报 — "看见趋势"\n面向技术研发与市场战略部门\n系统性追踪行业前沿动态\n\n[✓] 销售情报 — "抓住机会"\n面向销售和商务团队\n驱动业绩增长的实战利器',
       sz=11, c=MCK_DARK)
    # 可扩展领域（大卡片）
    tx(sl, Inches(5), Inches(2.5), Inches(4), Inches(0.4),
       '可快速扩展领域', sz=15, bold=True, c=MCK_BLUE_L)
    for i, (d, desc) in enumerate([
        ('[链接] 供应链管理', '寻源、评估、批准、监控，全链路可视化'),
        ('[文档] 知识产权监控', '专利追踪、侵权预警，保护核心资产'),
        ('[银行] 金融市场追踪', '投融资动态、并购重组信号，把握市场脉搏'),
        ('[人群] 人力资源情报', '竞对人事变动、人才流动趋势，优化组织策略'),
    ]):
        y = Inches(3.0 + i*0.7)
        tx(sl, Inches(5.2), y, Inches(4), Inches(0.6), f'{d} — {desc}', sz=11, c=MCK_DARK)
    # 效率
    rc(sl, Inches(1.5), Inches(5.0), Inches(7), Inches(0.8), bd=MCK_GREEN)
    tx(sl, Inches(1.5), Inches(5.1), Inches(7), Inches(0.6),
       '一份配置文件 + 一个前端模板 → 数天内上线新域\n所有通用逻辑都已收敛在共享内核中',
       sz=16, bold=True, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s12(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '安全合规：企业级安全内置', MCK_GREEN, 11)
    # 安全表格（大）
    rows = [
        ['安全维度', '实现方式', '安全等级'],
        ['认证', 'JWT Bearer Token（HS256），密码 SHA-256 + 随机盐存储', '企业级'],
        ['密钥保护', 'API Key / Agent Key 落地 XOR 混淆，API 响应中脱敏为 ***', '高'],
        ['CORS', '环境变量严格白名单控制允许来源域名', '高'],
        ['RBAC 权限', 'Admin / Manager / Analyst / Viewer 四级角色，精细到菜单级', '企业级'],
        ['审计日志', '所有变更操作记录操作人身份和时间戳，满足合规追溯', '企业级'],
    ]
    cw = [Inches(1.8), Inches(6.0), Inches(1.5)]
    mck_table(sl, Inches(0.5), Inches(1.2), Inches(9), rows, cw, hc=MCK_GREEN, ch=0.50)
    # 底部
    rc(sl, Inches(0.5), Inches(4.5), Inches(9), Inches(0.9), bd=MCK_GREEN)
    tx(sl, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.7),
       '每一个操作都被记录 · 每一个密钥都被保护 · 每一次访问都有据可查\n满足 SOC2、ISO27001 等合规要求',
       sz=15, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s13(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术优势：轻而不薄', MCK_BLUE_L, 12)
    # 对比表格（大）
    rows = [
        ['维度', '传统方案', 'Intelligence Web', '优势说明'],
        ['数据库', 'Oracle/PostgreSQL\n年费 10-50 万', 'SQLite 单文件\n零许可费', '节省 100% 数据库成本'],
        ['运维', '专职 DBA + DevOps\n年成本 30-60 万', '容器化一键部署\n无需专人运维', '省去 1 个全职岗位'],
        ['框架', 'React/Angular 企业版\nBI 工具许可', 'Vanilla JS + Flask\n全部开源免费', '节省 100% 框架成本'],
        ['部署周期', '数周至数月', 'docker compose up -d\n分钟级上线', '缩短 90%+ 部署时间'],
        ['扩展性', '改代码、重新部署\n测试', '配置文件 + 前端模板\n数天上线新域', '敏捷迭代，快速响应'],
    ]
    cw = [Inches(1.5), Inches(2.8), Inches(2.8), Inches(2.2)]
    mck_table(sl, Inches(0.3), Inches(1.2), Inches(9.4), rows, cw, hc=MCK_BLUE_L, ch=0.50)
    # 底部
    rc(sl, Inches(0.5), Inches(4.1), Inches(9), Inches(0.8), bd=MCK_BLUE_L)
    tx(sl, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.6),
       '你不需要再养一个团队来维护这套系统\n极简架构，让技术团队专注业务创新',
       sz=16, bold=True, c=MCK_BLUE_L, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s14(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '竞品对比：为什么不是 CRM？', MCK_RED, 13)
    # 对比表格（大）
    rows = [
        ['维度', '通用 CRM', 'OA 系统', 'Excel', 'Intelligence Web'],
        ['定位', '客户关系管理', '办公流程管理', '临时记录工具', '企业情报智能管理'],
        ['信息采集', '无自动采集功能', '无自动采集功能', '手动录入', 'AI 自动采集 + 增量同步'],
        ['AI 分析', '无 AI 分析能力', '无 AI 分析能力', '无 AI 分析能力', 'AI Agent 自动摘要分析'],
        ['多域扩展', '固定功能模块', '固定功能模块', '手动搭建', '配置文件数天上线新域'],
        ['部署成本', '百万级许可证', '十万级实施费', '低但无法规模化', '几乎为零（开源）'],
        ['情报深度', '浅（联系人 + 商机）', '浅（审批流程）', '无', '深（采集→分析→预警→行动）'],
    ]
    cw = [Inches(1.5), Inches(1.7), Inches(1.7), Inches(1.5), Inches(2.5)]
    mck_table(sl, Inches(0.2), Inches(1.2), Inches(9.6), rows, cw, hc=MCK_RED, ch=0.48)
    # 差异化
    rc(sl, Inches(0.5), Inches(5.0), Inches(9), Inches(1.0), bd=MCK_BLUE)
    bar(sl, Inches(0.5), Inches(5.0), Pt(5), Inches(1.0), MCK_BLUE)
    tx(sl, Inches(0.8), Inches(5.1), Inches(8.4), Inches(0.8),
       'Intelligence Web 只做一件事：帮企业把散落在各处的情报变成可行动的洞察。\n为此做了极深的垂直打磨，这是通用 CRM 和 OA 系统无法做到的。',
       sz=14, c=MCK_BLUE)
    pn(sl, 19)

def s15(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '适用人群：为谁而建？', MCK_BLUE, 14)
    users = [
        ('一线销售/商务经理', '比竞对更快知道客户在哪\n需求是什么', '自动预警 + 客户画像 + 商机全追踪', MCK_BLUE_L),
        ('市场研究/战略规划', '持续扫描行业全貌\n形成可指导决策的报告', '多渠道采集 + AI 分析 + 趋势可视化', MCK_GREEN),
        ('企业管理者/决策层', '一眼看清整体状况\n不做凭感觉的赌局', '数据看板 + AI 摘要 + 组织级公共资产', MCK_RED),
    ]
    for i, (role, need, sol, clr) in enumerate(users):
        l = Inches(0.2 + i*3.2)
        rc(sl, l, Inches(1.2), Inches(2.9), Inches(5.0), bd=clr)
        bar(sl, l+Inches(0.3), Inches(1.4), Inches(0.8), Pt(4), clr)
        tx(sl, l+Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.5), role, sz=15, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.5), Inches(2.5), Inches(1.0), need, sz=11, c=MCK_GRAY, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(3.7), Inches(2.5), Inches(0.3), '[箭头]', sz=16, c=clr, al=PP_ALIGN.CENTER)
        lb(sl, l+Inches(0.1), Inches(4.2), Inches(2.7), Inches(1.4), sol, clr, tc=WHITE, sz=11, bold=True)
    tx(sl, Inches(0.5), Inches(6.5), Inches(9), Inches(0.4),
       '不是给所有人的万能工具 — 是为情报驱动决策的团队量身定制的效率倍增器',
       sz=12, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s16(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '部署启动：最快数分钟', MCK_BLUE_L, 15)
    steps = [
        ('Step 1', 'docker compose up -d', '一条命令启动全部服务\nAPI + Nginx + Meilisearch 一键部署', MCK_BLUE_L),
        ('Step 2', '浏览器访问 :8765', '登录系统\n开始使用 15 个功能页面', MCK_GREEN),
        ('Step 3', '配置数据源 + 项目', '开始自动采集\nAI Agent 立即工作', MCK_BLUE),
    ]
    for i, (step, cmd, desc, clr) in enumerate(steps):
        l = Inches(0.5 + i*3.2)
        rc(sl, l, Inches(1.2), Inches(2.9), Inches(2.8), bd=clr)
        bar(sl, l+Inches(0.3), Inches(1.4), Inches(0.8), Pt(4), clr)
        tx(sl, l+Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.5), step, sz=18, bold=True, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.4), Inches(2.5), Inches(0.5), cmd, sz=12, c=MCK_BLUE_L, al=PP_ALIGN.CENTER, bold=True)
        tx(sl, l+Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.7), desc, sz=10, c=MCK_GRAY, al=PP_ALIGN.CENTER)
        if i < 2:
            tx(sl, l+Inches(2.85), Inches(2.3), Inches(0.4), Inches(0.4), '→', sz=20, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    # 开箱即用
    tx(sl, Inches(0.5), Inches(4.4), Inches(4), Inches(0.4),
       '开箱即用的能力', sz=16, bold=True, c=MCK_BLUE_L)
    feats = [
        '[✓] 15 个功能页面，即开即用，无需配置',
        '[✓] RBAC 权限体系，三级角色即刻生效',
        '[✓] AI Agent 预设模板，配置即用，无需训练',
        '[✓] 暗色模式、响应式设计，现代用户体验',
        '[✓] 新业务域：一份配置 + 一个模板 → 数天内上线',
    ]
    for i, f in enumerate(feats):
        tx(sl, Inches(0.5), Inches(4.9 + i*0.45), Inches(9), Inches(0.4),
           f, sz=12, c=MCK_DARK)
    pn(sl, 19)

def s17(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '投资回报：1-3 个月回本', MCK_GREEN, 16)
    # 三个大数字
    nums = [
        ('1-3', '个月\n投资回报周期', MCK_GREEN),
        ('50-100', '万\n年新增收入', MCK_BLUE_L),
        ('0.5-1', '个\n全职人力释放', MCK_BLUE),
    ]
    for i, (num, sfx, clr) in enumerate(nums):
        x = Inches(0.5 + i*3.2)
        rc(sl, x, Inches(1.2), Inches(2.9), Inches(2.5), bd=clr)
        bar(sl, x+Inches(0.3), Inches(1.4), Inches(0.8), Pt(4), clr)
        bn(sl, x+Inches(0.1), Inches(1.8), Inches(2.7), Inches(1.3), num, sfx=sfx, c=clr, fs=40)
    # 基准
    rc(sl, Inches(0.5), Inches(4.0), Inches(9), Inches(0.6))
    tx(sl, Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.4),
       '基准假设：一名年薪资 20 万的销售人员，平均每月处理 50+ 条商机线索', sz=12, c=MCK_GRAY)
    # 详细数据
    data = [
        ('人力释放', 'AI 替代 1-2 小时/天 → 相当于 0.5-1 个全职人力释放', MCK_BLUE_L),
        ('新增商机收入', '转化率提升 20% + 销售周期缩短 20% → 额外 50-100 万年收入', MCK_GREEN),
        ('系统部署成本', '几乎为零（开源 + 自有服务器即可运行）', MCK_GRAY),
        ('投资回报周期', '1-3 个月即可收回全部投资成本', MCK_RED),
    ]
    for i, (k, v, clr) in enumerate(data):
        y = Inches(4.8 + i*0.5)
        rc(sl, Inches(0.5), y, Inches(9), Inches(0.45), bd=clr)
        tx(sl, Inches(0.7), y+Inches(0.08), Inches(2.2), Inches(0.28), k, sz=12, bold=True, c=clr)
        tx(sl, Inches(3.0), y+Inches(0.08), Inches(6.3), Inches(0.28), v, sz=10, c=MCK_DARK)
    pn(sl, 19)

def s18(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '行动号召：升级您的企业情报管理', MCK_BLUE, 17)
    # 核心信息（大卡片）
    pts = [
        ('[网络]', '开源 · 免费 · 可私有化部署', '零成本启动，数据完全自主可控'),
        ('[AI]', 'AI 驱动 · 人机协同 · 持续进化', '系统越用越聪明，形成正向飞轮'),
        ('[图表]', '1-3 个月投资回报 · 100 万年新增收入', '高杠杆投入，快速收回成本'),
        ('[火箭]', '一条命令启动 · 数天扩展新域', '极简部署，敏捷迭代'),
    ]
    for i, (icon, title, desc) in enumerate(pts):
        y = Inches(1.3 + i*0.65)
        rc(sl, Inches(1.5), y, Inches(7), Inches(0.55), bd=BORD)
        tx(sl, Inches(1.7), y+Inches(0.08), Inches(1.5), Inches(0.38), icon, sz=12, c=MCK_BLUE, al=PP_ALIGN.CENTER)
        tx(sl, Inches(3.3), y+Inches(0.08), Inches(2.5), Inches(0.38), title, sz=13, bold=True, c=MCK_DARK)
        tx(sl, Inches(5.9), y+Inches(0.08), Inches(2.4), Inches(0.38), desc, sz=10, c=MCK_GRAY)
    # 联系方式
    rc(sl, Inches(2), Inches(4.3), Inches(6), Inches(1.8), bd=MCK_GREEN)
    tx(sl, Inches(2.2), Inches(4.5), Inches(5.6), Inches(0.6),
       '准备好让您的企业情报管理升级了吗？',
       sz=20, bold=True, c=MCK_GREEN, al=PP_ALIGN.CENTER)
    tx(sl, Inches(2.2), Inches(5.3), Inches(5.6), Inches(0.4),
       '[官网]  [邮箱]  [电话]  [二维码]',
       sz=12, c=MCK_GRAY, al=PP_ALIGN.CENTER)
    pn(sl, 19)

def s19(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG)
    bar(sl, Inches(0), Inches(0), Inches(10), Pt(5), MCK_BLUE)
    tx(sl, Inches(1), Inches(2.5), Inches(8), Inches(1.0),
       'Intelligence Web', sz=48, bold=True, c=MCK_DARK, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.6), Inches(8), Inches(0.6),
       '让情报成为您的核心竞争力', sz=22, c=MCK_BLUE_L, al=PP_ALIGN.CENTER)
    bar(sl, Inches(4.3), Inches(4.4), Inches(1.4), Pt(3), MCK_GREEN)
    tx(sl, Inches(2), Inches(4.8), Inches(6), Inches(1.0),
       '谢谢\nThank You',
       sz=24, c=MCK_LIGHT, al=PP_ALIGN.CENTER)
    pn(sl, 19)

# ============================================================
# 主函数
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s1(prs); s2(prs); s3(prs); s4(prs); s5(prs); s6(prs); s8(prs); s9(prs)
    s7(prs); s10(prs); s11(prs); s12(prs); s13(prs); s14(prs); s15(prs)
    s16(prs); s17(prs); s18(prs); s19(prs)

    prs.save('docs/Intelligence_Web_麦肯锡风格 PPT_v2.pptx')
    print('PPT 麦肯锡风格 v2 已生成：docs/Intelligence_Web_麦肯锡风格 PPT_v2.pptx')
    print('19 页 · 16:9 宽屏 · 内容充实版 · 大字体 · 满内容 · 含架构图')

if __name__ == '__main__':
    main()
#!/usr/bin/env python3
"""
Intelligence Web — 产品宣发 PPT 生成器 v6
浅色商务主题：白色背景 · 黑体 · 专业简洁 · 16:9 宽屏
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# 浅色商务配色
# ============================================================
BG = RGBColor(0xF8, 0xF9, 0xFA)         # 浅灰背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)        # 纯白
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)      # 卡片白
BORD = RGBColor(0xE2, 0xE8, 0xF0)         # 浅灰边框

PRI = RGBColor(0x1A, 0x56, 0xDB)          # 主蓝
PRI_L = RGBColor(0x4C, 0x84, 0xF5)        # 蓝亮
ACC = RGBColor(0x0F, 0x9D, 0x58)          # 翠绿
ACC_L = RGBColor(0x34, 0xD3, 0x99)        # 绿亮
WRN = RGBColor(0xFA, 0xBD, 0x24)          # 琥珀
WRN_D = RGBColor(0xF5, 0x9E, 0x0B)        # 琥珀深
PUR = RGBColor(0x7C, 0x3A, 0xED)          # 紫
PUR_L = RGBColor(0xA7, 0x8B, 0xF6)        # 紫亮
PNK = RGBColor(0xEC, 0x48, 0x99)          # 粉
CYN = RGBColor(0x08, 0x91, 0xB2)          # 青
TXT = RGBColor(0x1F, 0x29, 0x37)          # 主文字（深灰）
SEC = RGBColor(0x6B, 0x72, 0x80)          # 次要文字
MUT = RGBColor(0x9C, 0xA1, 0xAB)          # 弱化文字
DIV = RGBColor(0xE5, 0xE7, 0xEB)          # 分隔线

FONT_CN = 'SimHei'    # 黑体
FONT_NUM = 'Arial'    # 数字

# ============================================================
# 工具函数
# ============================================================

def set_bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def shadow(s, blur=Pt(12), dist=Pt(4)):
    try:
        sp = s.element.spPr
        el = sp.get_or_add_extLst(); xml = el.get_or_add_extLst()
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        o = etree.SubElement(xml, f'{{{ns}}}outerShdw')
        o.set('blurRad', str(blur.emu)); o.set('dist', str(dist.emu))
        o.set('dir', '2700000'); o.set('algn', 'ctr')
        sc = etree.SubElement(o, f'{{{ns}}}srClr')
        a = etree.SubElement(sc, f'{{{ns}}}a'); a.set('val', '30000')
        so = etree.SubElement(sc, f'{{{ns}}}solidClr')
        sr = etree.SubElement(so, f'{{{ns}}}srgbClr'); sr.set('val', '000000')
    except: pass

def card(slide, l, t, w, h, bg=CARD_BG, bd=BORD, r=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = r
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.color.rgb = bd; s.line.width = Pt(0.5)
    shadow(s)
    return s

def bar(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def circ(slide, l, t, sz, c, op=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    try:
        sf = s.fill._fill; n = sf.find(qn('a:solidFill'))
        if n is not None and len(n) > 0:
            al = etree.SubElement(n[0], qn('a:alpha')); al.set('val', str(int(op*1000)))
    except: pass
    return s

def tx(slide, l, t, w, h, text, sz=18, bold=False, c=TXT, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = FONT_CN
    p.alignment = al; p.space_after = Pt(2)
    return tb

def ml(slide, l, t, w, h, *lines, sz=15, c=TXT, al=PP_ALIGN.LEFT,
       bul=False, sp=1.5, bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    flat = []
    for item in lines:
        if isinstance(item, list): flat.extend(item)
        else: flat.append(item)
    for i, line in enumerate(flat):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ('▸ ' if bul and line.strip() else '') + line
        p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = FONT_CN
        p.alignment = al; p.space_after = Pt(6)
        if sp != 1.0: p.line_spacing = Pt(sz * sp)
    return tb

def bn(slide, l, t, w, h, num, sfx='', c=ACC, fs=52):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(fs); p.font.bold = True; p.font.color.rgb = c; p.font.name = FONT_NUM
    p.alignment = PP_ALIGN.CENTER
    if sfx:
        p2 = tf.add_paragraph(); p2.text = sfx
        p2.font.size = Pt(15); p2.font.color.rgb = SEC; p2.font.name = FONT_CN
        p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    return tb

def lb(slide, l, t, w, h, text, bg, tc=WHITE, sz=12, bold=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.35
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = tc; p.font.bold = bold; p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    return s

def pn(slide, n):
    tb = slide.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.6), Pt(10))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = f'{n}/18'; p.font.size = Pt(9); p.font.color.rgb = MUT; p.font.name = FONT_NUM
    p.alignment = PP_ALIGN.RIGHT

def hdr(slide, title, c, n):
    tx(slide, Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.6),
       title, sz=26, bold=True, c=TXT)
    bar(slide, Inches(0.6), Inches(0.85), Inches(0.8), Pt(3), c)
    pn(slide, n)

# ============================================================
# 卡片式表格
# ============================================================

def ctbl(slide, l, t, w, rows, cw, hc=PRI, ch=0.5, gap=0):
    """传统表格（恢复为 pptx 原生表格）"""
    ncols = len(rows[0])
    nrows = len(rows)
    ts = slide.shapes.add_table(nrows, ncols, l, t, w, ch * nrows + gap * (nrows-1))
    table = ts.table
    # 设置列宽
    cx = 0
    for i, col_w in enumerate(cw):
        table.columns[i].width = col_w
    # 填充数据
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.text = rows[r][c] if r < len(rows) and c < len(rows[r]) else ""
            # 设置单元格格式
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9 if r > 0 else 10)
                para.font.name = FONT_CN
                para.font.color.rgb = WHITE if r == 0 else TXT
                para.font.bold = r == 0
                para.alignment = PP_ALIGN.LEFT
            # 设置单元格背景色
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hc
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return ts

def kv(slide, l, t, w, h, k, v, kc=SEC, vc=TXT):
    """键值卡片"""
    hw = (w - Inches(0.05))/2
    card(slide, l, t, hw, h, bg=CARD_BG, bd=BORD, r=0.03)
    tx(slide, l+Inches(0.06), t+Inches(0.06), hw-Inches(0.12), h-Inches(0.12),
       k, sz=10, c=kc)
    card(slide, l+hw+Inches(0.05), t, hw, h, bg=RGBColor(0xF8, 0xF9, 0xFA), bd=BORD, r=0.03)
    tx(slide, l+hw+Inches(0.05)+Inches(0.06), t+Inches(0.06), hw-Inches(0.12), h-Inches(0.12),
       v, sz=9, c=vc)

# ============================================================
# 背景装饰
# ============================================================

def bg_d(slide):
    circ(slide, Inches(8.5), Inches(-1), Inches(4), PRI, 0.05)
    circ(slide, Inches(-0.8), Inches(6.5), Inches(2), ACC, 0.04)

def cover_d(slide):
    circ(slide, Inches(-1), Inches(-1.5), Inches(6), PRI, 0.06)
    circ(slide, Inches(7.5), Inches(4), Inches(5), PUR, 0.04)
    bar(slide, Inches(2.5), Inches(1.2), Inches(5), Pt(2), PRI)
    bar(slide, Inches(3.5), Inches(6.5), Inches(3), Pt(2), ACC)

def close_d(slide):
    circ(slide, Inches(2.5), Inches(0.5), Inches(7), PRI, 0.04)
    circ(slide, Inches(6), Inches(4.5), Inches(5), PUR, 0.03)
    bar(slide, Inches(3), Inches(1.3), Inches(4), Pt(2), PRI)
    bar(slide, Inches(4), Inches(7), Inches(2), Pt(2), ACC)

# ============================================================
# 各页
# ============================================================

def s1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); cover_d(sl)
    # 顶部色条
    bar(sl, Inches(0), Inches(0), Inches(10), Pt(4), PRI)
    tx(sl, Inches(1), Inches(2.2), Inches(8), Inches(1),
       'Intelligence Web', sz=48, bold=True, c=TXT, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.4), Inches(8), Inches(0.6),
       '企业情报智能管理平台', sz=22, c=PRI, al=PP_ALIGN.CENTER)
    bar(sl, Inches(4.2), Inches(4.2), Inches(1.6), Pt(2), ACC)
    tx(sl, Inches(1.5), Inches(4.5), Inches(7), Inches(0.9),
       '从情报采集到行动闭环\n让每一个决策都有据可依',
       sz=15, c=SEC, al=PP_ALIGN.CENTER)
    for i, tag in enumerate(['开源免费', 'AI 驱动', '容器化部署']):
        lb(sl, Inches(3+i*1.8), Inches(5.9), Inches(1.5), Inches(0.4),
           tag, CARD_BG, tc=PRI, sz=11, bold=True)
    pn(sl, 1)

def s2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '我们正处在一个信息就是竞争力的时代', WRN_D, 2)
    ml(sl, Inches(0.6), Inches(1.2), Inches(8.2), Inches(4),
       '全球商业环境变化速度加快，竞争对手的动作从"月级"缩短到"天级"',
       '企业每年在信息采集上投入大量人力，但 90% 的信息被淹没、遗忘或被竞争对手捷足先登',
       '传统的信息搜集方式：碎片化、被动式、依赖个人经验',
       '市场空白：缺乏一套系统化的平台，把"散落的线索"变成"可行动的洞察"',
       sz=14, bul=True, sp=1.6)
    card(sl, Inches(6.5), Inches(1.8), Inches(2.9), Inches(2.2))
    tx(sl, Inches(6.7), Inches(1.95), Inches(2.5), Inches(0.4),
       '关键数据', sz=12, bold=True, c=WRN_D)
    bn(sl, Inches(6.7), Inches(2.45), Inches(2.5), Inches(1),
       '90%', sfx='信息被淹没或遗忘', c=WRN_D, fs=44)
    pn(sl, 2)

def s3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '企业情报管理的三大顽疾', WRN_D, 3)
    pains = [
        ('01', '信息孤岛', '各部门各管一套，\n重要情报反复丢失', '错失机会 · 重复劳动', WRN_D),
        ('02', '被动应对', '竞争对手出手后才\n反应过来', '永远慢半拍', PNK),
        ('03', '决策靠直觉', '管理层拍板靠经验\n和感觉', '高风险 · 事后才知', PUR),
    ]
    for i, (num, title, desc, cost, clr) in enumerate(pains):
        l = Inches(0.3 + i*3.2)
        card(sl, l, Inches(1.3), Inches(2.9), Inches(5.2), bd=clr)
        bar(sl, l+Inches(0.3), Inches(1.5), Inches(0.8), Pt(3), clr)
        tx(sl, l, Inches(1.8), Inches(2.9), Inches(0.5), num, sz=32, bold=True, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l, Inches(2.4), Inches(2.9), Inches(0.45), title, sz=18, bold=True, c=TXT, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(3), Inches(2.5), Inches(1.3), desc, sz=11, c=SEC, al=PP_ALIGN.CENTER)
        bar(sl, l+Inches(0.7), Inches(4.6), Inches(1.5), Pt(1), BORD)
        tx(sl, l+Inches(0.2), Inches(4.8), Inches(2.5), Inches(0.6), cost, sz=11, c=clr, al=PP_ALIGN.CENTER)
    pn(sl, 3)

def s4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, 'Intelligence Web 是什么？', PRI, 4)
    card(sl, Inches(0.6), Inches(1.3), Inches(8.8), Inches(1.6), bd=PRI)
    bar(sl, Inches(0.6), Inches(1.3), Pt(4), Inches(1.6), PRI)
    tx(sl, Inches(0.9), Inches(1.5), Inches(8.2), Inches(1.2),
       '情报采集 → 结构化存储 → AI 分析 → 行动闭环\n一体化平台',
       sz=22, bold=True, c=TXT, al=PP_ALIGN.CENTER)
    extras = [
        ('[搜索]', '系统追踪行业动态、竞争对手动向、客户需求变化', '和潜在商业机会'),
        ('[大脑]', '不再依赖碎片化信息搜集', '和个人经验判断'),
        ('[图表]', '让每一个决策都有据可依', '数据驱动，而非直觉'),
    ]
    for i, (icon, l1, l2) in enumerate(extras):
        l = Inches(0.3 + i*3.2)
        card(sl, l, Inches(3.3), Inches(2.9), Inches(1.5))
        tx(sl, l+Inches(0.2), Inches(3.5), Inches(2.5), Inches(0.4), icon, sz=12, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(3.95), Inches(2.5), Inches(0.45), l1, sz=11, c=TXT)
        tx(sl, l+Inches(0.2), Inches(4.4), Inches(2.5), Inches(0.35), l2, sz=11, c=SEC)
    flow = ['[采集]', '[存储]', '[分析]', '[行动]']
    for i, item in enumerate(flow):
        x = Inches(2.1 + i*1.6)
        lb(sl, x, Inches(5.3), Inches(1.3), Inches(0.45), item, PRI, tc=WHITE, sz=11, bold=True)
        if i < len(flow)-1:
            tx(sl, x+Inches(1.3), Inches(5.32), Inches(0.3), Inches(0.35),
               '->', sz=13, c=MUT, al=PP_ALIGN.CENTER)
    pn(sl, 4)

def s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '五大能力，覆盖情报管理全链路', CYN, 5)
    feats = [
        ('[AI]', 'AI 智能分析', 'AI Agent 自动阅读、摘要、分析、留批注\n人机协同知识积累', PUR),
        ('[搜索]', '多渠道采集', '网站抓取、API 对接\n按日/周/月灵活设定采集频率', PRI),
        ('[图表]', '数据看板', '仪表盘统计、多维筛选排序\nMeilisearch 毫秒级全文搜索', CYN),
        ('[目标]', '商机管理', '从线索到成交的全生命周期追踪\n销售域专属', ACC),
        ('[齿轮]', '系统管控', 'RBAC 四级权限、审计日志\n通知中心、个性化配置', WRN_D),
    ]
    pos = [(0.3, 1.3), (3.5, 1.3), (6.7, 1.3), (1.9, 4.1), (5.1, 4.1)]
    for i, (icon, title, desc, clr) in enumerate(feats):
        l, tp = pos[i]
        card(sl, Inches(l), Inches(tp), Inches(2.9), Inches(2.6))
        bar(sl, Inches(l)+Inches(0.3), Inches(tp)+Inches(0.1), Inches(0.8), Pt(3), clr)
        tx(sl, Inches(l)+Inches(0.2), Inches(tp)+Inches(0.3), Inches(2.5), Inches(0.4),
           icon, sz=11, al=PP_ALIGN.CENTER)
        tx(sl, Inches(l)+Inches(0.2), Inches(tp)+Inches(0.75), Inches(2.5), Inches(0.4),
           title, sz=14, bold=True, c=TXT)
        tx(sl, Inches(l)+Inches(0.2), Inches(tp)+Inches(1.25), Inches(2.5), Inches(0.7),
           desc, sz=10, c=SEC)
    pn(sl, 5)

def s6(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '制造情报域 — "看见趋势"', ACC, 6)
    tx(sl, Inches(0.6), Inches(1.2), Inches(4.2), Inches(0.35),
       '核心数据流', sz=14, bold=True, c=ACC)
    flows = [
        '竞品发布新产品 → AI 5 分钟内捕获、摘要、推送预警',
        '国家发布智能制造补贴政策 → 自动抓取、分析、标记关联度',
        '行业论坛讨论下一代工艺 → 实时追踪、归类、形成趋势报告',
    ]
    for i, f in enumerate(flows):
        y = Inches(1.8 + i*1)
        card(sl, Inches(0.6), y, Inches(4.1), Inches(0.8))
        bar(sl, Inches(0.6), y+Inches(0.25), Pt(3), Inches(0.25), ACC)
        tx(sl, Inches(0.9), y+Inches(0.12), Inches(3.6), Inches(0.5), f, sz=11, c=TXT)
    tx(sl, Inches(5.3), Inches(1.2), Inches(3.7), Inches(0.35),
       '价值体现', sz=14, bold=True, c=ACC)
    vals = ['实时监控竞品产品线调整和产能扩张计划',
            '预判下一代制造技术的商业化时间表',
            '从"事后追悔"到"事前预判"']
    for i, v in enumerate(vals):
        y = Inches(1.8 + i*1)
        card(sl, Inches(5.3), y, Inches(3.6), Inches(0.8))
        tx(sl, Inches(5.5), y+Inches(0.2), Inches(3.2), Inches(0.45),
           f'[ok] {v}', sz=11, c=TXT)
    card(sl, Inches(1), Inches(5.4), Inches(8), Inches(0.7), bd=ACC)
    tx(sl, Inches(1), Inches(5.45), Inches(8), Inches(0.6),
       '看见趋势 · 看见未来 · 看见机会', sz=18, bold=True, c=ACC, al=PP_ALIGN.CENTER)
    pn(sl, 6)

def s7(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '销售情报域 — "抓住机会"', ACC, 7)
    scns = [
        ('[目标] 客户扩产预警', '客户宣布扩产 → 第一时间触发跟进流程', PRI),
        ('[工厂] 竞对投资信号', '竞争对手新建工厂前就被捕捉到投资信号', WRN_D),
        ('[趋势] 商机漏斗可视化', '销售主管通过数据看板掌握团队商机漏斗健康度', ACC),
    ]
    for i, (icon, desc, clr) in enumerate(scns):
        l = Inches(0.3 + i*3.2)
        card(sl, l, Inches(1.3), Inches(2.9), Inches(2))
        bar(sl, l+Inches(0.3), Inches(1.5), Inches(0.8), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.75), Inches(2.5), Inches(0.4), icon, sz=11, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.2), Inches(2.5), Inches(0.4), desc, sz=10, c=TXT)
    tx(sl, Inches(0.6), Inches(3.7), Inches(5), Inches(0.3),
       '商机全生命周期', sz=13, bold=True, c=SEC)
    stages = ['待核实', '合格商机', '方案报价', '商务谈判', '成交/丢标']
    scs = [MUT, PRI, CYN, ACC, WRN_D]
    for i, (st, sc) in enumerate(zip(stages, scs)):
        x = Inches(0.15 + i*1.9)
        active = i < 4
        lb(sl, x, Inches(4.2), Inches(1.7), Inches(0.45), st,
           sc if active else CARD_BG, tc=TXT if not active else WHITE, sz=10, bold=True)
        if i < len(stages)-1:
            tx(sl, x+Inches(1.7), Inches(4.18), Inches(0.3), Inches(0.35),
               '->', sz=12, c=MUT, al=PP_ALIGN.CENTER)
    pn(sl, 7)

def s8(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '"共享内核 + 任意域扩展"', PUR, 8)
    card(sl, Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.7), bd=PUR)
    bar(sl, Inches(0.6), Inches(1.3), Pt(4), Inches(0.7), PUR)
    tx(sl, Inches(0.9), Inches(1.4), Inches(8.2), Inches(0.5),
       '平台不设固定的业务边界，客户需要什么领域，就搭建什么领域',
       sz=14, c=PUR)
    tx(sl, Inches(0.6), Inches(2.4), Inches(3), Inches(0.3),
       '已验证领域', sz=13, bold=True, c=ACC)
    card(sl, Inches(0.6), Inches(2.8), Inches(3.6), Inches(1.3))
    tx(sl, Inches(0.8), Inches(2.95), Inches(3.2), Inches(0.9),
       '[ok] 制造情报\n"看见趋势"\n\n[ok] 销售情报\n"抓住机会"',
       sz=12, c=TXT)
    tx(sl, Inches(4.9), Inches(2.4), Inches(4), Inches(0.3),
       '可快速扩展领域', sz=13, bold=True, c=PRI)
    for i, (d, desc) in enumerate([
        ('[链接] 供应链管理', '寻源、评估、批准、监控'),
        ('[文档] 知识产权监控', '专利追踪、侵权预警'),
        ('[银行] 金融市场追踪', '投融资动态、并购重组'),
        ('[人群] 人力资源情报', '竞对人事变动、人才流动'),
    ]):
        y = Inches(2.8 + i*0.65)
        tx(sl, Inches(5.1), y, Inches(4), Inches(0.5), f'{d} — {desc}', sz=11, c=TXT)
    card(sl, Inches(1.5), Inches(5.4), Inches(7), Inches(0.7), bd=ACC)
    tx(sl, Inches(1.5), Inches(5.45), Inches(7), Inches(0.6),
       '一份配置文件 + 一个前端模板 → 数天内上线新域',
       sz=16, bold=True, c=ACC, al=PP_ALIGN.CENTER)
    pn(sl, 8)

def s9(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '为 IT 负责人准备的架构透明度', CYN, 9)
    rows = [
        ['服务', '端口', '技术栈', '角色'],
        ['Research API', '8766', 'Flask + Python 3.11', '制造情报业务引擎'],
        ['Sales API', '8767', 'Gunicorn + Python 3.11', '销售情报业务引擎'],
        ['Gateway', '8765', 'Nginx Alpine', '反向代理 + JWT 鉴权'],
        ['Meilisearch', '7700', 'Meilisearch v1.12', '全文检索引擎'],
    ]
    cw = [Inches(1.8), Inches(0.7), Inches(2.6), Inches(2.8)]
    ctbl(sl, Inches(0.5), Inches(1.3), Inches(9), rows, cw, hc=CYN, ch=0.5)
    tx(sl, Inches(0.6), Inches(4.2), Inches(5), Inches(0.3),
       '关键架构原则', sz=14, bold=True, c=CYN)
    principles = [
        '[共享] 共享内核 — 所有业务域通过 Docker Volume 挂载同一 core/ 目录',
        '[分离] 分离领域 — 各域数据完全隔离，通过 SQLite 独立文件实现',
        '[接口] MCP Server — 18 个工具方法，AI Agent 通过标准化协议直接访问数据',
        '[轻量] 零外部数据库依赖 — 除 Meilisearch 外无需任何额外服务',
    ]
    for i, p in enumerate(principles):
        tx(sl, Inches(0.6), Inches(4.6 + i*0.5), Inches(8.5), Inches(0.4),
           p, sz=11, c=TXT)
    pn(sl, 9)

def s10(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '"轻"而不"薄"的技术底座', ACC, 10)
    rows = [
        ['维度', '传统方案', 'Intelligence Web'],
        ['数据库', 'Oracle/PostgreSQL\n年费 10-50 万', 'SQLite 单文件\n零许可费'],
        ['运维', '专职 DBA + DevOps\n年成本 30-60 万', '容器化一键部署\n无需专人运维'],
        ['框架', 'React/Angular 企业版\nBI 工具许可', 'Vanilla JS + Flask\n全部开源免费'],
        ['部署周期', '数周至数月', 'docker compose up -d\n分钟级上线'],
        ['扩展性', '改代码、重新部署\n测试', '配置文件 + 前端模板\n数天上线新域'],
    ]
    cw = [Inches(1.8), Inches(3.5), Inches(3.5)]
    ctbl(sl, Inches(0.3), Inches(1.3), Inches(9.4), rows, cw, hc=ACC, ch=0.55)
    card(sl, Inches(1), Inches(5.6), Inches(8), Inches(0.6), bd=ACC)
    tx(sl, Inches(1), Inches(5.65), Inches(8), Inches(0.5),
       '你不需要再养一个团队来维护这套系统', sz=16, bold=True, c=ACC, al=PP_ALIGN.CENTER)
    pn(sl, 10)

def s11(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '不是工具，是协作伙伴', PUR, 11)
    steps = [
        ('[AI] AI Agent\n每日自动巡检\n指定数据源', PRI),
        ('[分析] AI 对采集内容\n进行初步分析\n和归类', CYN),
        ('[建议] AI 在每条情报下\n留下观察和\n建议', PUR),
        ('[决策] 人类员工\n做出最终判断\n和行动决策', ACC),
        ('[反馈] 人类反馈\n反哺 AI，持续\n提升精度', WRN_D),
    ]
    for i, (text, clr) in enumerate(steps):
        x = Inches(0.1 + i*1.95)
        card(sl, x, Inches(1.3), Inches(1.75), Inches(1.4), bd=clr)
        bar(sl, x+Inches(0.2), Inches(1.5), Inches(1.35), Pt(2), clr)
        tx(sl, x+Inches(0.1), Inches(1.75), Inches(1.55), Inches(0.85),
           text, sz=10, c=TXT, al=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            tx(sl, x+Inches(1.65), Inches(1.75), Inches(0.35), Inches(0.35),
               '->', sz=14, bold=True, c=clr, al=PP_ALIGN.CENTER)
    card(sl, Inches(0.6), Inches(3.2), Inches(8.8), Inches(3), bd=PUR)
    pts = [
        '这不是一个"录入工具"，是让人类和 AI Agent 共同工作的平台',
        '每一次人工反馈都在训练 AI，形成正向飞轮',
        '18 个 MCP 工具方法，支持 Claude / Hermes / OpenClaw 等多 Agent 接入',
    ]
    for i, p in enumerate(pts):
        tx(sl, Inches(0.9), Inches(3.5 + i*0.6), Inches(8.2), Inches(0.45),
           f'[dot] {p}', sz=12, c=TXT)
    pn(sl, 11)

def s12(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '企业级安全，从第一天就内置', ACC, 12)
    rows = [
        ['安全维度', '实现方式'],
        ['认证', 'JWT Bearer Token（HS256），密码 SHA-256 + 随机盐存储'],
        ['密钥保护', 'API Key / Agent Key 落地 XOR 混淆，API 响应中脱敏为 ***'],
        ['CORS', '环境变量严格白名单控制允许来源域名'],
        ['RBAC 权限', 'Admin / Manager / Analyst / Viewer 四级角色，精细到菜单级'],
        ['审计日志', '所有变更操作记录操作人身份和时间戳，满足合规追溯'],
    ]
    cw = [Inches(1.8), Inches(6.2)]
    ctbl(sl, Inches(0.7), Inches(1.3), Inches(8.6), rows, cw, hc=ACC, ch=0.5)
    card(sl, Inches(1), Inches(5.3), Inches(8), Inches(0.7), bd=ACC)
    tx(sl, Inches(1), Inches(5.35), Inches(8), Inches(0.6),
       '每一个操作都被记录 · 每一个密钥都被保护 · 每一次访问都有据可查',
       sz=13, c=ACC, al=PP_ALIGN.CENTER)
    pn(sl, 12)

def s13(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '每一分钱，都能算得清楚', ACC, 13)
    tx(sl, Inches(0.6), Inches(1.2), Inches(4), Inches(0.3),
       '成本对比（年）', sz=14, bold=True, c=ACC)
    rows = [
        ['项目', '传统方案', 'Intelligence Web', '节省'],
        ['数据库许可', '10-50 万/年', '0 元', '100%'],
        ['运维人力', '30-60 万/年', '0 元', '省去 1 个全职'],
        ['框架/工具许可', '10-30 万/年', '0 元', '100%'],
        ['部署周期', '数周至数月', '分钟级', '90%+ 时间'],
    ]
    cw = [Inches(1.8), Inches(2.1), Inches(2.7), Inches(1.7)]
    ctbl(sl, Inches(0.3), Inches(1.6), Inches(9.4), rows, cw, hc=ACC, ch=0.5)
    tx(sl, Inches(0.6), Inches(4.5), Inches(5), Inches(0.3),
       '效率提升', sz=14, bold=True, c=ACC)
    effs = [('10x', '信息采集'), ('50x', '情报分析'), ('100x', '信息检索'), ('10x', '商机响应')]
    for i, (num, label) in enumerate(effs):
        x = Inches(0.3 + i*2.4)
        bn(sl, x, Inches(4.9), Inches(2.2), Inches(1.1), num, sfx=label, c=ACC, fs=40)
    pn(sl, 13)

def s14(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '以中型企业销售团队为例，算一笔账', PRI, 14)
    card(sl, Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.5))
    tx(sl, Inches(0.8), Inches(1.35), Inches(8.4), Inches(0.4),
       '基准假设：一名年薪资 20 万的销售人员', sz=12, c=SEC)
    nums = [
        ('1-3 个月', '投资回报周期', ACC),
        ('50-100 万', '年新增收入', PRI),
        ('0.5-1 个', '全职人力释放', PUR),
    ]
    for i, (num, label, clr) in enumerate(nums):
        x = Inches(0.5 + i*3.2)
        card(sl, x, Inches(2.1), Inches(2.9), Inches(2), bd=clr)
        bar(sl, x+Inches(0.3), Inches(2.3), Inches(0.8), Pt(3), clr)
        bn(sl, x+Inches(0.1), Inches(2.6), Inches(2.7), Inches(1),
           num, sfx=label, c=clr, fs=36)
    data = [
        ('人力释放', 'AI 替代 1-2 小时/天 → 相当于 0.5-1 个全职人力释放'),
        ('新增商机收入', '转化率提升 20% + 销售周期缩短 20% → 额外 50-100 万年收入'),
        ('系统部署成本', '几乎为零（开源 + 自有服务器）'),
        ('投资回报周期', '1-3 个月'),
    ]
    for i, (k, v) in enumerate(data):
        y = Inches(4.5 + i*0.5)
        kv(sl, Inches(0.6), y, Inches(8.8), Inches(0.42), k, v)
    tx(sl, Inches(1), Inches(6.6), Inches(8), Inches(0.4),
       '这不是一个成本中心，这是一个收入引擎',
       sz=18, bold=True, c=ACC, al=PP_ALIGN.CENTER)
    pn(sl, 14)

def s15(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '为谁而建？', PRI, 15)
    users = [
        ('一线销售/商务经理', '比竞对更快知道客户在哪\n需求是什么',
         '自动预警 + 客户画像 + 商机全追踪', PRI),
        ('市场研究/战略规划', '持续扫描行业全貌\n形成可指导决策的报告',
         '多渠道采集 + AI 分析 + 趋势可视化', CYN),
        ('企业管理者/决策层', '一眼看清整体状况\n不做凭感觉的赌局',
         '数据看板 + AI 摘要 + 组织级公共资产', ACC),
    ]
    for i, (role, need, sol, clr) in enumerate(users):
        l = Inches(0.2 + i*3.2)
        card(sl, l, Inches(1.3), Inches(2.9), Inches(5.4), bd=clr)
        bar(sl, l+Inches(0.3), Inches(1.5), Inches(0.8), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.4),
           role, sz=14, bold=True, c=TXT, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.3), Inches(2.5), Inches(1),
           need, sz=10, c=SEC, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(3.6), Inches(2.5), Inches(0.3),
           '[箭头]', sz=16, bold=True, c=clr, al=PP_ALIGN.CENTER)
        lb(sl, l+Inches(0.1), Inches(4.1), Inches(2.7), Inches(1.2),
           sol, clr, tc=WHITE, sz=9, bold=True)
    tx(sl, Inches(1), Inches(6.9), Inches(8), Inches(0.3),
       '不是给所有人的万能工具 — 是为情报驱动决策的团队量身定制的效率倍增器',
       sz=11, c=MUT, al=PP_ALIGN.CENTER)
    pn(sl, 15)

def s16(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '为什么不是 CRM？不是 OA？不是 Excel？', WRN_D, 16)
    rows = [
        ['维度', '通用 CRM', 'OA 系统', 'Excel', 'Intelligence Web'],
        ['定位', '客户关系管理', '办公流程管理', '临时记录', '企业情报智能管理'],
        ['信息采集', '[x]', '[x]', '[x]', 'AI 自动采集'],
        ['AI 分析', '[x]', '[x]', '[x]', 'AI Agent 分析'],
        ['多域扩展', '固定模块', '固定模块', '手动搭建', '数天上线新域'],
        ['部署成本', '百万级', '十万级', '低', '几乎为零'],
        ['情报深度', '浅', '浅', '无', '深（全链路）'],
    ]
    cw = [Inches(1.5), Inches(1.7), Inches(1.7), Inches(1.5), Inches(2.2)]
    ctbl(sl, Inches(0.2), Inches(1.3), Inches(9.6), rows, cw, hc=WRN_D, ch=0.45)
    card(sl, Inches(0.6), Inches(5.3), Inches(8.8), Inches(0.9), bd=PRI)
    bar(sl, Inches(0.6), Inches(5.3), Pt(4), Inches(0.9), PRI)
    tx(sl, Inches(0.9), Inches(5.4), Inches(8.2), Inches(0.7),
       'Intelligence Web 只做一件事：帮企业把散落在各处的情报变成可行动的洞察。为此做了极深的垂直打磨。',
       sz=12, c=PRI)
    pn(sl, 16)

def s17(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); bg_d(sl); hdr(sl, '最快数分钟，启动您的企业情报系统', PRI, 17)
    steps = [
        ('Step 1', 'docker compose up -d', '一条命令启动全部服务', PRI),
        ('Step 2', '浏览器访问 :8765', '登录，开始使用', ACC),
        ('Step 3', '配置数据源 + 项目', '开始自动采集', CYN),
    ]
    for i, (step, cmd, desc, clr) in enumerate(steps):
        l = Inches(0.5 + i*3.2)
        card(sl, l, Inches(1.3), Inches(2.9), Inches(2.4), bd=clr)
        bar(sl, l+Inches(0.3), Inches(1.5), Inches(0.8), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.4),
           step, sz=16, bold=True, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.25), Inches(2.5), Inches(0.45),
           cmd, sz=11, c=PRI, al=PP_ALIGN.CENTER, bold=True)
        tx(sl, l+Inches(0.2), Inches(2.8), Inches(2.5), Inches(0.4),
           desc, sz=11, c=SEC, al=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            tx(sl, l+Inches(2.85), Inches(2.2), Inches(0.4), Inches(0.4),
               '->', sz=18, bold=True, c=MUT, al=PP_ALIGN.CENTER)
    tx(sl, Inches(0.6), Inches(4.1), Inches(4), Inches(0.3),
       '开箱即用的能力', sz=14, bold=True, c=PRI)
    feats = [
        '[ok] 15 个功能页面，即开即用',
        '[ok] RBAC 权限体系，三级角色即刻生效',
        '[ok] AI Agent 预设模板，配置即用',
        '[ok] 暗色模式、响应式设计，现代用户体验',
        '[ok] 新业务域：一份配置 + 一个模板 → 数天内上线',
    ]
    for i, f in enumerate(feats):
        tx(sl, Inches(0.6), Inches(4.5 + i*0.45), Inches(8.5), Inches(0.35),
           f, sz=11, c=TXT)
    pn(sl, 17)

def s18(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); close_d(sl)
    bar(sl, Inches(0), Inches(0), Inches(10), Pt(4), PRI)
    tx(sl, Inches(1), Inches(2), Inches(8), Inches(0.9),
       'Intelligence Web', sz=44, bold=True, c=TXT, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.1), Inches(8), Inches(0.5),
       '让情报成为您的核心竞争力', sz=20, c=PRI, al=PP_ALIGN.CENTER)
    bar(sl, Inches(4.2), Inches(3.8), Inches(1.6), Pt(2), ACC)
    pts = [
        '[网络] 开源 · 免费 · 可私有化部署',
        '[AI] AI 驱动 · 人机协同 · 持续进化',
        '[图表] 1-3 个月投资回报 · 100 万年新增收入',
        '[火箭] 一条命令启动 · 数天扩展新域',
    ]
    for i, p in enumerate(pts):
        tx(sl, Inches(2), Inches(4.1 + i*0.4), Inches(6), Inches(0.3),
           p, sz=13, c=TXT, al=PP_ALIGN.CENTER)
    card(sl, Inches(2), Inches(6), Inches(6), Inches(0.6), bd=ACC)
    tx(sl, Inches(2), Inches(6.05), Inches(6), Inches(0.5),
       '准备好让您的企业情报管理升级了吗？',
       sz=18, bold=True, c=ACC, al=PP_ALIGN.CENTER)
    tx(sl, Inches(2), Inches(6.85), Inches(6), Inches(0.2),
       '[官网]  [邮箱]  [电话]  [二维码]', sz=9, c=MUT, al=PP_ALIGN.CENTER)
    pn(sl, 18)

# ============================================================
# 主函数
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s1(prs); s2(prs); s3(prs); s4(prs); s5(prs); s6(prs); s7(prs); s8(prs)
    s9(prs); s10(prs); s11(prs); s12(prs); s13(prs); s14(prs); s15(prs)
    s16(prs); s17(prs); s18(prs)

    prs.save('docs/Intelligence_Web_产品宣发 PPT.pptx')
    print('PPT v6 已生成：docs/Intelligence_Web_产品宣发 PPT.pptx')
    print('18 页 · 16:9 宽屏 · 黑体 · 浅色商务主题')

if __name__ == '__main__':
    main()
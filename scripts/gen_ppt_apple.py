#!/usr/bin/env python3
"""
Intelligence Web — 苹果风格 PPT 生成器
极简设计 · 大量留白 · 精致布局 · 16:9 宽屏
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ============================================================
# 苹果配色
# ============================================================
BG = RGBColor(0xFF, 0xFF, 0xFF)         # 纯白
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_100 = RGBColor(0xF5, 0xF5, 0xF5)   # 浅灰
GRAY_200 = RGBColor(0xE5, 0xE5, 0xE5)   # 边框灰
GRAY_400 = RGBColor(0xA0, 0xA0, 0xA0)   # 次要文字
GRAY_600 = RGBColor(0x6E, 0x6E, 0x6E)   # 正文文字
GRAY_800 = RGBColor(0x33, 0x33, 0x33)   # 主文字
GRAY_900 = RGBColor(0x1D, 0x1D, 0x1F)   # 标题文字

APPLE_BLUE = RGBColor(0x00, 0x71, 0xE3)  # 苹果蓝
APPLE_GREEN = RGBColor(0x34, 0xC7, 0x59)  # 苹果绿
APPLE_RED = RGBColor(0xFF, 0x3B, 0x30)    # 苹果红
APPLE_ORANGE = RGBColor(0xFF, 0x95, 0x00) # 苹果橙

FONT_CN = 'PingFang SC'    # 苹方（苹果风格）
FONT_EN = 'SF Pro Display' # 苹果官方字体（备选 Arial）

# ============================================================
# 工具函数（苹果风格）
# ============================================================

def set_bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def add_shadow(s, blur=Pt(20), dist=Pt(8)):
    """添加苹果风格阴影"""
    try:
        sp = s.element.spPr
        el = sp.get_or_add_extLst(); xml = el.get_or_add_extLst()
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        o = etree.SubElement(xml, f'{{{ns}}}outerShdw')
        o.set('blurRad', str(blur.emu)); o.set('dist', str(dist.emu))
        o.set('dir', '2700000'); o.set('algn', 'ctr')
        sc = etree.SubElement(o, f'{{{ns}}}srClr')
        a = etree.SubElement(sc, f'{{{ns}}}a'); a.set('val', '15000')
        so = etree.SubElement(sc, f'{{{ns}}}solidClr')
        sr = etree.SubElement(so, f'{{{ns}}}srgbClr'); sr.set('val', '000000')
    except: pass

def apple_card(slide, l, t, w, h, bg=WHITE, bd=GRAY_200, r=0.08):
    """苹果风格卡片"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = r
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.color.rgb = bd; s.line.width = Pt(0.5)
    add_shadow(s)
    return s

def apple_bar(slide, l, t, w, h, c):
    """装饰条"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def tx(slide, l, t, w, h, text, sz=24, bold=False, c=GRAY_800, al=PP_ALIGN.LEFT):
    """文本框"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = FONT_CN
    p.alignment = al; p.space_after = Pt(4)
    return tb

def bn(slide, l, t, w, h, num, sfx='', c=APPLE_BLUE, fs=64):
    """大数字"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(fs); p.font.bold = True; p.font.color.rgb = c; p.font.name = FONT_EN
    p.alignment = PP_ALIGN.CENTER
    if sfx:
        p2 = tf.add_paragraph(); p2.text = sfx
        p2.font.size = Pt(18); p2.font.color.rgb = GRAY_400; p2.font.name = FONT_CN
        p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)
    return tb

def lb(slide, l, t, w, h, text, bg=APPLE_BLUE, tc=WHITE, sz=12, bold=False):
    """标签"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = tc; p.font.bold = bold; p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    return s

def pn(slide, n, total=18):
    """页码"""
    tb = slide.shapes.add_textbox(Inches(9.3), Inches(7.15), Inches(0.6), Pt(10))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = f'{n}/{total}'; p.font.size = Pt(9); p.font.color.rgb = GRAY_400; p.font.name = FONT_EN
    p.alignment = PP_ALIGN.RIGHT

def hdr(slide, title, c=APPLE_BLUE, n=1):
    """标题"""
    tx(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7),
       title, sz=32, bold=True, c=GRAY_900)
    pn(slide, n)

def add_image(slide, path, l, t, w, h):
    """添加图片"""
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
        return True
    return False

# ============================================================
# 各页（苹果风格）
# ============================================================

def s1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG)
    # 顶部装饰线
    apple_bar(sl, Inches(0), Inches(0), Inches(10), Pt(3), APPLE_BLUE)
    # 主标题
    tx(sl, Inches(1), Inches(2.5), Inches(8), Inches(1.2),
       'Intelligence Web', sz=56, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.8), Inches(8), Inches(0.7),
       '企业情报智能管理平台', sz=24, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    apple_bar(sl, Inches(4.3), Inches(4.7), Inches(1.4), Pt(2), APPLE_BLUE)
    tx(sl, Inches(1.5), Inches(5.0), Inches(7), Inches(1),
       '从情报采集到行动闭环\n让每一个决策都有据可依',
       sz=16, c=GRAY_600, al=PP_ALIGN.CENTER)
    # 标签
    for i, tag in enumerate(['开源免费', 'AI 驱动', '容器化部署']):
        lb(sl, Inches(3+i*1.8), Inches(6.3), Inches(1.5), Inches(0.4),
           tag, WHITE, tc=APPLE_BLUE, sz=11, bold=True)
    pn(sl, 1)

def s2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '执行摘要', APPLE_BLUE, 2)
    # 核心观点
    apple_card(sl, Inches(0.6), Inches(1.3), Inches(8.8), Inches(1.4))
    tx(sl, Inches(0.9), Inches(1.5), Inches(8.2), Inches(1.1),
       '企业每年在信息采集上投入大量人力，但 90% 的信息被淹没、遗忘或被竞争对手捷足先登',
       sz=22, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
    # 三个痛点
    pts = [
        ('信息孤岛', '各部门各管一套\n重要情报反复丢失', '错失机会', APPLE_BLUE),
        ('被动应对', '竞争对手出手后才\n反应过来', '永远慢半拍', APPLE_RED),
        ('决策靠直觉', '管理层拍板靠经验\n和感觉', '高风险赌局', GRAY_600),
    ]
    for i, (title, desc, cost, clr) in enumerate(pts):
        l = Inches(0.4 + i*3.2)
        apple_card(sl, l, Inches(3.1), Inches(2.8), Inches(3.2))
        apple_bar(sl, l+Inches(0.3), Inches(3.3), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(3.6), Inches(2.4), Inches(0.5), title, sz=18, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(4.2), Inches(2.4), Inches(0.8), desc, sz=12, c=GRAY_600, al=PP_ALIGN.CENTER)
        apple_bar(sl, l+Inches(0.7), Inches(5.3), Inches(1.4), Pt(1), GRAY_200)
        tx(sl, l+Inches(0.2), Inches(5.5), Inches(2.4), Inches(0.4), cost, sz=11, c=clr, al=PP_ALIGN.CENTER)
    pn(sl, 2)

def s3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '解决方案', APPLE_BLUE, 3)
    # 核心定义
    apple_card(sl, Inches(0.6), Inches(1.3), Inches(8.8), Inches(1.8), bd=APPLE_BLUE)
    apple_bar(sl, Inches(0.6), Inches(1.3), Pt(4), Inches(1.8), APPLE_BLUE)
    tx(sl, Inches(0.9), Inches(1.5), Inches(8.2), Inches(1.5),
       '情报采集 → 结构化存储 → AI 分析 → 行动闭环\n一体化平台',
       sz=26, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
    # 四个支柱
    pillars = [
        ('AI 智能分析', '自动阅读、摘要、分析、留批注\n人机协同知识积累', APPLE_BLUE),
        ('多渠道采集', '网站抓取、API 对接\n按日/周/月灵活设定频率', APPLE_GREEN),
        ('数据看板', '仪表盘统计、多维筛选\n毫秒级全文搜索', APPLE_BLUE),
        ('商机管理', '从线索到成交的全生命周期追踪\n销售域专属', APPLE_RED),
    ]
    for i, (title, desc, clr) in enumerate(pillars):
        l = Inches(0.3 + i*2.5)
        apple_card(sl, l, Inches(3.5), Inches(2.3), Inches(2.6))
        apple_bar(sl, l+Inches(0.25), Inches(3.7), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(4.0), Inches(1.9), Inches(0.5), title, sz=14, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(4.6), Inches(1.9), Inches(0.8), desc, sz=11, c=GRAY_600, al=PP_ALIGN.CENTER)
    pn(sl, 3)

def s4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '核心价值', APPLE_GREEN, 4)
    # 四个大数字
    effs = [
        ('10x', '信息采集', 'AI 自动采集替代人工\n每日 1-2 小时解放', APPLE_BLUE),
        ('50x', '情报分析', 'AI 秒级完成人工\n15 分钟摘要工作', APPLE_GREEN),
        ('100x', '信息检索', '毫秒级全文搜索\n中英文混合检索', APPLE_BLUE),
        ('10x', '商机响应', '分钟级预警推送\n第一时间触发跟进', APPLE_RED),
    ]
    for i, (num, title, desc, clr) in enumerate(effs):
        x = Inches(0.3 + i*2.5)
        apple_card(sl, x, Inches(1.3), Inches(2.3), Inches(2.4))
        bn(sl, x+Inches(0.1), Inches(1.5), Inches(2.1), Inches(1.0), num, c=clr, fs=56)
        tx(sl, x+Inches(0.1), Inches(2.7), Inches(2.1), Inches(0.4), title, sz=14, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
        tx(sl, x+Inches(0.1), Inches(3.2), Inches(2.1), Inches(0.5), desc, sz=10, c=GRAY_600, al=PP_ALIGN.CENTER)
    # 底部结论
    apple_card(sl, Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.9), bd=APPLE_GREEN)
    tx(sl, Inches(0.9), Inches(4.3), Inches(8.2), Inches(0.7),
       'Intelligence Web 不是成本中心，而是收入引擎\n预计为中型企业带来额外 50-100 万年收入',
       sz=18, bold=True, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 4)

def s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '投资回报', APPLE_BLUE, 5)
    # 成本对比表格
    rows = [
        ['项目', '传统方案', 'Intelligence Web', '节省'],
        ['数据库许可', '年费 10-50 万', '零许可费', '100%'],
        ['运维人力', '年成本 30-60 万', '无需专人运维', '省去 1 个全职'],
        ['框架许可', '企业版许可', '全部开源免费', '100%'],
        ['部署周期', '数周至数月', '分钟级上线', '90%+ 时间'],
    ]
    # 简单表格布局
    y_start = Inches(1.3)
    for r, row in enumerate(rows):
        y = y_start + r * Inches(0.55)
        if r == 0:
            apple_card(sl, Inches(0.6), y, Inches(8.8), Inches(0.45), bd=APPLE_BLUE)
            for c, txt in enumerate(row):
                w = Inches(2.0) if c < 3 else Inches(1.5)
                tx(sl, Inches(0.7+c*2.15), y+Inches(0.08), w, Inches(0.3), txt, sz=11, bold=True, c=WHITE)
        else:
            bg = WHITE if r%2==0 else GRAY_100
            apple_card(sl, Inches(0.6), y, Inches(8.8), Inches(0.45), bd=GRAY_200)
            for c, txt in enumerate(row):
                w = Inches(2.0) if c < 3 else Inches(1.5)
                tc = APPLE_GREEN if '0' in txt and '节省' not in txt else GRAY_800
                tx(sl, Inches(0.7+c*2.15), y+Inches(0.08), w, Inches(0.3), txt, sz=10, c=tc)
    # 收益
    tx(sl, Inches(0.6), Inches(3.8), Inches(5), Inches(0.4),
       '收益量化（中型企业销售团队）', sz=16, bold=True, c=GRAY_900)
    data = [
        ('人力释放', 'AI 替代 1-2 小时/天 → 相当于 0.5-1 个全职'),
        ('新增收入', '转化率提升 20% + 周期缩短 20% → 额外 50-100 万年收入'),
        ('部署成本', '几乎为零（开源 + 自有服务器）'),
        ('回报周期', '1-3 个月'),
    ]
    for i, (k, v) in enumerate(data):
        y = Inches(4.3 + i*0.5)
        apple_card(sl, Inches(0.6), y, Inches(4.5), Inches(0.42), bd=GRAY_200)
        tx(sl, Inches(0.8), y+Inches(0.08), Inches(1.4), Inches(0.25), k, sz=11, bold=True, c=APPLE_BLUE)
        tx(sl, Inches(2.3), y+Inches(0.08), Inches(2.6), Inches(0.25), v, sz=9, c=GRAY_800)
    bn(sl, Inches(6.5), Inches(4.0), Inches(3), Inches(1.4),
       '1-3', sfx='个月\n投资回报周期', c=APPLE_GREEN, fs=48)
    pn(sl, 5)

def s6(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术架构', APPLE_BLUE, 6)
    # 架构图
    # 客户端层
    apple_card(sl, Inches(3.2), Inches(1.3), Inches(3.6), Inches(0.7), bd=APPLE_BLUE)
    tx(sl, Inches(3.2), Inches(1.4), Inches(3.6), Inches(0.5),
       '客户端\n浏览器 Portal SPA', sz=10, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    tx(sl, Inches(4.8), Inches(2.05), Inches(0.4), Inches(0.3), '↓', sz=14, c=GRAY_400, al=PP_ALIGN.CENTER)
    # 网关层
    apple_card(sl, Inches(2.7), Inches(2.4), Inches(4.6), Inches(0.7), bd=APPLE_GREEN)
    tx(sl, Inches(2.7), Inches(2.5), Inches(4.6), Inches(0.5),
       '网关层\nNginx 反向代理 + JWT 鉴权', sz=10, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    tx(sl, Inches(4.8), Inches(3.15), Inches(0.4), Inches(0.3), '↓', sz=14, c=GRAY_400, al=PP_ALIGN.CENTER)
    # 业务域层
    apple_card(sl, Inches(0.6), Inches(3.6), Inches(4.2), Inches(0.7), bd=APPLE_BLUE)
    tx(sl, Inches(0.6), Inches(3.7), Inches(4.2), Inches(0.5),
       '制造情报 API\nFlask + Python 3.11', sz=10, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    apple_card(sl, Inches(5.2), Inches(3.6), Inches(4.2), Inches(0.7), bd=APPLE_GREEN)
    tx(sl, Inches(5.2), Inches(3.7), Inches(4.2), Inches(0.5),
       '销售情报 API\nGunicorn + Python 3.11', sz=10, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    tx(sl, Inches(2.5), Inches(4.35), Inches(0.4), Inches(0.3), '↓', sz=14, c=GRAY_400, al=PP_ALIGN.CENTER)
    tx(sl, Inches(7.5), Inches(4.35), Inches(0.4), Inches(0.3), '↓', sz=14, c=GRAY_400, al=PP_ALIGN.CENTER)
    # 核心引擎层
    apple_card(sl, Inches(1.8), Inches(4.7), Inches(6.4), Inches(0.7), bd=APPLE_BLUE)
    tx(sl, Inches(1.8), Inches(4.8), Inches(6.4), Inches(0.5),
       '共享核心引擎 core/\napp.py · db.py · datasource.py · project.py', sz=10, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    # 数据层
    apple_card(sl, Inches(0.6), Inches(5.7), Inches(2.5), Inches(0.6), bd=APPLE_GREEN)
    tx(sl, Inches(0.6), Inches(5.8), Inches(2.5), Inches(0.4), 'SQLite\nresearch.db', sz=9, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    apple_card(sl, Inches(3.5), Inches(5.7), Inches(2.5), Inches(0.6), bd=APPLE_GREEN)
    tx(sl, Inches(3.5), Inches(5.8), Inches(2.5), Inches(0.4), 'SQLite\nsales.db', sz=9, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    apple_card(sl, Inches(6.4), Inches(5.7), Inches(2.5), Inches(0.6), bd=APPLE_BLUE)
    tx(sl, Inches(6.4), Inches(5.8), Inches(2.5), Inches(0.4), 'Meilisearch\n全文检索', sz=9, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    pn(sl, 6)

def s7(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术栈', APPLE_BLUE, 7)
    # 技术栈列表
    techs = [
        ('Flask', 'Python 3.11', '制造情报业务引擎'),
        ('Gunicorn', 'Python 3.11', '销售情报业务引擎'),
        ('Nginx', 'Alpine Linux', '反向代理 + JWT 鉴权'),
        ('Meilisearch', 'v1.12', '全文检索引擎'),
        ('SQLite', '单文件', '零配置数据库'),
        ('Vanilla JS', 'HTML/CSS/JS', '15 个功能页面'),
        ('MCP Server', 'JSON-RPC', '18 个工具方法'),
    ]
    for i, (name, version, desc) in enumerate(techs):
        y = Inches(1.3 + i*0.65)
        apple_card(sl, Inches(0.6), y, Inches(8.8), Inches(0.55))
        tx(sl, Inches(0.8), y+Inches(0.1), Inches(2.0), Inches(0.35), name, sz=13, bold=True, c=GRAY_900)
        tx(sl, Inches(2.9), y+Inches(0.1), Inches(1.5), Inches(0.35), version, sz=11, c=GRAY_600)
        tx(sl, Inches(4.5), y+Inches(0.1), Inches(4.5), Inches(0.35), desc, sz=10, c=GRAY_600)
    # 底部
    apple_card(sl, Inches(0.6), Inches(6.2), Inches(8.8), Inches(0.7), bd=APPLE_BLUE)
    tx(sl, Inches(0.9), Inches(6.3), Inches(8.2), Inches(0.5),
       '极简技术栈，经过充分验证，稳定可靠且易于维护\n无需 DBA，无需 DevOps，一条命令即可启动整套系统',
       sz=13, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    pn(sl, 7)

def s8(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, 'AI 协作', APPLE_GREEN, 8)
    # 五步流程
    steps = [
        ('AI 自动巡检', '每日自动扫描\n指定数据源', APPLE_BLUE),
        ('AI 分析归类', '对采集内容\n进行初步分析', APPLE_GREEN),
        ('AI 留下建议', '在每条情报下\n留下观察建议', APPLE_BLUE),
        ('人类决策', '做出最终判断\n和行动决策', APPLE_RED),
        ('反馈训练 AI', '人类反馈\n反哺 AI 提升', APPLE_BLUE),
    ]
    for i, (title, desc, clr) in enumerate(steps):
        x = Inches(0.15 + i*1.95)
        apple_card(sl, x, Inches(1.3), Inches(1.7), Inches(1.5), bd=clr)
        apple_bar(sl, x+Inches(0.2), Inches(1.5), Inches(1.3), Pt(2), clr)
        tx(sl, x+Inches(0.1), Inches(1.7), Inches(1.5), Inches(0.4), title, sz=11, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
        tx(sl, x+Inches(0.1), Inches(2.2), Inches(1.5), Inches(0.5), desc, sz=9, c=GRAY_600, al=PP_ALIGN.CENTER)
        if i < 4:
            tx(sl, x+Inches(1.7), Inches(1.9), Inches(0.3), Inches(0.3), '→', sz=14, c=GRAY_400, al=PP_ALIGN.CENTER)
    # 核心观点
    apple_card(sl, Inches(0.6), Inches(3.2), Inches(8.8), Inches(1.4), bd=APPLE_GREEN)
    tx(sl, Inches(0.9), Inches(3.4), Inches(8.2), Inches(1.1),
       '这不是一个"录入工具"，是让人类和 AI Agent 共同工作的平台\n每一次人工反馈都在训练 AI，形成正向飞轮效应',
       sz=16, c=GRAY_900, al=PP_ALIGN.CENTER)
    pn(sl, 8)

def s9(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '制造情报', APPLE_GREEN, 9)
    # 左侧：数据流
    tx(sl, Inches(0.6), Inches(1.3), Inches(4.5), Inches(0.4),
       '核心数据流', sz=16, bold=True, c=APPLE_GREEN)
    flows = [
        '竞品发布新产品 → AI 5 分钟内捕获、摘要、推送预警',
        '国家发布智能制造补贴政策 → 自动抓取、分析、标记关联度',
        '行业论坛讨论下一代工艺 → 实时追踪、归类、形成趋势报告',
    ]
    for i, f in enumerate(flows):
        y = Inches(1.9 + i*1.0)
        apple_card(sl, Inches(0.6), y, Inches(4.2), Inches(0.85))
        apple_bar(sl, Inches(0.6), y+Inches(0.25), Pt(3), Inches(0.25), APPLE_GREEN)
        tx(sl, Inches(0.9), y+Inches(0.15), Inches(3.7), Inches(0.55), f, sz=11, c=GRAY_800)
    # 右侧：价值
    tx(sl, Inches(5.4), Inches(1.3), Inches(4), Inches(0.4),
       '价值体现', sz=16, bold=True, c=APPLE_GREEN)
    vals = ['实时监控竞品产品线调整和产能扩张计划',
            '预判下一代制造技术的商业化时间表',
            '从"事后追悔"到"事前预判"']
    for i, v in enumerate(vals):
        y = Inches(1.9 + i*1.0)
        apple_card(sl, Inches(5.4), y, Inches(4.0), Inches(0.85))
        tx(sl, Inches(5.6), y+Inches(0.2), Inches(3.6), Inches(0.5),
           f'✓ {v}', sz=11, c=GRAY_800)
    # 底部
    apple_card(sl, Inches(0.6), Inches(5.3), Inches(8.8), Inches(0.7), bd=APPLE_GREEN)
    tx(sl, Inches(0.9), Inches(5.4), Inches(8.2), Inches(0.5),
       '看见趋势 · 看见未来 · 看见机会',
       sz=20, bold=True, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 9)

def s10(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '销售情报', APPLE_BLUE, 10)
    # 三个场景
    scns = [
        ('客户扩产预警', '客户宣布扩产 → 第一时间触发跟进流程', APPLE_BLUE),
        ('竞对投资信号', '竞争对手新建工厂前就被捕捉到投资信号', APPLE_RED),
        ('商机漏斗可视化', '销售主管通过数据看板掌握团队商机漏斗健康度', APPLE_GREEN),
    ]
    for i, (title, desc, clr) in enumerate(scns):
        l = Inches(0.3 + i*3.2)
        apple_card(sl, l, Inches(1.3), Inches(2.9), Inches(2.2))
        apple_bar(sl, l+Inches(0.3), Inches(1.5), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.5), title, sz=14, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.4), Inches(2.5), Inches(0.6), desc, sz=10, c=GRAY_600)
    # 商机漏斗
    tx(sl, Inches(0.6), Inches(3.9), Inches(5), Inches(0.4),
       '商机全生命周期', sz=16, bold=True, c=GRAY_900)
    stages = ['待核实', '合格商机', '方案报价', '商务谈判', '成交/丢标']
    scs = [GRAY_400, APPLE_BLUE, APPLE_GREEN, APPLE_BLUE, APPLE_RED]
    for i, (st, sc) in enumerate(zip(stages, scs)):
        x = Inches(0.2 + i*1.9)
        lb(sl, x, Inches(4.4), Inches(1.7), Inches(0.45), st, sc, tc=WHITE if i<4 else GRAY_800, sz=11, bold=True)
        if i < 4:
            tx(sl, x+Inches(1.7), Inches(4.35), Inches(0.3), Inches(0.35), '→', sz=12, c=GRAY_400, al=PP_ALIGN.CENTER)
    pn(sl, 10)

def s11(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '扩展性', APPLE_BLUE, 11)
    # 核心理念
    apple_card(sl, Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.8), bd=APPLE_BLUE)
    apple_bar(sl, Inches(0.6), Inches(1.3), Pt(4), Inches(0.8), APPLE_BLUE)
    tx(sl, Inches(0.9), Inches(1.4), Inches(8.2), Inches(0.6),
       '平台不设固定的业务边界，客户需要什么领域，就搭建什么领域\n从提出需求到新域可用，最快数天内完成交付',
       sz=14, c=GRAY_900)
    # 已验证
    tx(sl, Inches(0.6), Inches(2.5), Inches(3), Inches(0.4),
       '已验证领域', sz=14, bold=True, c=APPLE_GREEN)
    apple_card(sl, Inches(0.6), Inches(3.0), Inches(3.8), Inches(1.4))
    tx(sl, Inches(0.8), Inches(3.15), Inches(3.4), Inches(1.1),
       '✓ 制造情报 — "看见趋势"\n面向技术研发与市场战略部门\n系统性追踪行业前沿动态',
       sz=11, c=GRAY_800)
    tx(sl, Inches(5.0), Inches(2.5), Inches(4), Inches(0.4),
       '可快速扩展领域', sz=14, bold=True, c=APPLE_BLUE)
    for i, (d, desc) in enumerate([
        ('供应链管理', '寻源、评估、批准、监控'),
        ('知识产权监控', '专利追踪、侵权预警'),
        ('金融市场追踪', '投融资动态、并购重组'),
        ('人力资源情报', '竞对人事变动、人才流动'),
    ]):
        y = Inches(3.0 + i*0.65)
        tx(sl, Inches(5.2), y, Inches(4), Inches(0.55), f'{d} — {desc}', sz=11, c=GRAY_800)
    # 效率
    apple_card(sl, Inches(1.5), Inches(5.2), Inches(7), Inches(0.7), bd=APPLE_GREEN)
    tx(sl, Inches(1.5), Inches(5.3), Inches(7), Inches(0.5),
       '一份配置文件 + 一个前端模板 → 数天内上线新域\n所有通用逻辑都已收敛在共享内核中',
       sz=15, bold=True, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 11)

def s12(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '安全合规', APPLE_GREEN, 12)
    # 安全表格
    rows = [
        ['认证', 'JWT Bearer Token（HS256），密码 SHA-256 + 随机盐存储'],
        ['密钥保护', 'API Key / Agent Key 落地 XOR 混淆，API 响应中脱敏'],
        ['CORS', '环境变量严格白名单控制允许来源域名'],
        ['RBAC 权限', 'Admin / Manager / Analyst / Viewer 四级角色'],
        ['审计日志', '所有变更操作记录操作人身份和时间戳'],
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.3 + i*0.65)
        apple_card(sl, Inches(0.6), y, Inches(8.8), Inches(0.55))
        tx(sl, Inches(0.8), y+Inches(0.1), Inches(2.0), Inches(0.35), k, sz=12, bold=True, c=APPLE_GREEN)
        tx(sl, Inches(3.0), y+Inches(0.1), Inches(6.2), Inches(0.35), v, sz=10, c=GRAY_800)
    # 底部
    apple_card(sl, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.8), bd=APPLE_GREEN)
    tx(sl, Inches(0.9), Inches(4.9), Inches(8.2), Inches(0.6),
       '每一个操作都被记录 · 每一个密钥都被保护 · 每一次访问都有据可查\n满足 SOC2、ISO27001 等合规要求',
       sz=13, c=APPLE_GREEN, al=PP_ALIGN.CENTER)
    pn(sl, 12)

def s13(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '技术优势', APPLE_BLUE, 13)
    # 对比
    rows = [
        ['数据库', '年费 10-50 万', 'SQLite 单文件\n零许可费', '节省 100%'],
        ['运维', '年成本 30-60 万', '容器化一键部署\n无需专人运维', '省去 1 个全职'],
        ['框架', '企业版许可', 'Vanilla JS + Flask\n全部开源免费', '节省 100%'],
        ['部署', '数周至数月', 'docker compose up -d\n分钟级上线', '缩短 90%+'],
        ['扩展', '改代码、重新部署', '配置文件 + 前端模板\n数天上线新域', '敏捷迭代'],
    ]
    for i, (dim, trad, intel, adv) in enumerate(rows):
        y = Inches(1.3 + i*0.65)
        apple_card(sl, Inches(0.6), y, Inches(8.8), Inches(0.55))
        tx(sl, Inches(0.8), y+Inches(0.1), Inches(1.2), Inches(0.35), dim, sz=11, bold=True, c=GRAY_900)
        tx(sl, Inches(2.1), y+Inches(0.1), Inches(2.2), Inches(0.35), trad, sz=9, c=GRAY_600)
        tx(sl, Inches(4.4), y+Inches(0.1), Inches(2.2), Inches(0.35), intel, sz=9, c=APPLE_BLUE)
        tx(sl, Inches(6.7), y+Inches(0.1), Inches(2.5), Inches(0.35), adv, sz=9, c=APPLE_GREEN)
    # 底部
    apple_card(sl, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.7), bd=APPLE_BLUE)
    tx(sl, Inches(0.9), Inches(4.9), Inches(8.2), Inches(0.5),
       '你不需要再养一个团队来维护这套系统\n极简架构，让技术团队专注业务创新',
       sz=15, bold=True, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    pn(sl, 13)

def s14(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '竞品对比', APPLE_RED, 14)
    # 对比表格
    rows = [
        ['定位', '客户关系管理', '办公流程管理', '临时记录', '企业情报智能管理'],
        ['信息采集', '无', '无', '手动录入', 'AI 自动采集'],
        ['AI 分析', '无', '无', '无', 'AI Agent 分析'],
        ['多域扩展', '固定模块', '固定模块', '手动搭建', '数天上线新域'],
        ['部署成本', '百万级', '十万级', '低', '几乎为零'],
        ['情报深度', '浅', '浅', '无', '深（全链路）'],
    ]
    for i, row in enumerate(rows):
        y = Inches(1.3 + i*0.55)
        if i == 0:
            apple_card(sl, Inches(0.4), y, Inches(9.2), Inches(0.45), bd=APPLE_RED)
            for c, txt in enumerate(row):
                w = Inches(1.6) if c < 4 else Inches(2.0)
                tx(sl, Inches(0.5+c*1.7), y+Inches(0.08), w, Inches(0.3), txt, sz=10, bold=True, c=WHITE)
        else:
            bg = WHITE if i%2==0 else GRAY_100
            apple_card(sl, Inches(0.4), y, Inches(9.2), Inches(0.45), bd=GRAY_200)
            for c, txt in enumerate(row):
                w = Inches(1.6) if c < 4 else Inches(2.0)
                tc = APPLE_GREEN if c == 4 and 'AI' in txt else GRAY_800
                tx(sl, Inches(0.5+c*1.7), y+Inches(0.08), w, Inches(0.3), txt, sz=9, c=tc)
    # 差异化
    apple_card(sl, Inches(0.6), Inches(4.9), Inches(8.8), Inches(0.9), bd=APPLE_BLUE)
    apple_bar(sl, Inches(0.6), Inches(4.9), Pt(4), Inches(0.9), APPLE_BLUE)
    tx(sl, Inches(0.9), Inches(5.0), Inches(8.2), Inches(0.7),
       'Intelligence Web 只做一件事：帮企业把散落在各处的情报变成可行动的洞察\n为此做了极深的垂直打磨，这是通用 CRM 和 OA 系统无法做到的',
       sz=12, c=APPLE_BLUE)
    pn(sl, 14)

def s15(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '适用人群', APPLE_BLUE, 15)
    users = [
        ('一线销售/商务经理', '比竞对更快知道客户在哪\n需求是什么', '自动预警 + 客户画像 + 商机全追踪', APPLE_BLUE),
        ('市场研究/战略规划', '持续扫描行业全貌\n形成可指导决策的报告', '多渠道采集 + AI 分析 + 趋势可视化', APPLE_GREEN),
        ('企业管理者/决策层', '一眼看清整体状况\n不做凭感觉的赌局', '数据看板 + AI 摘要 + 组织级公共资产', APPLE_RED),
    ]
    for i, (role, need, sol, clr) in enumerate(users):
        l = Inches(0.2 + i*3.2)
        apple_card(sl, l, Inches(1.3), Inches(2.9), Inches(5.0), bd=clr)
        apple_bar(sl, l+Inches(0.3), Inches(1.5), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.5), role, sz=14, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.6), Inches(2.5), Inches(0.9), need, sz=11, c=GRAY_600, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(3.8), Inches(2.5), Inches(0.3), '↓', sz=16, c=clr, al=PP_ALIGN.CENTER)
        lb(sl, l+Inches(0.1), Inches(4.3), Inches(2.7), Inches(1.3), sol, clr, tc=WHITE, sz=10, bold=True)
    tx(sl, Inches(0.5), Inches(6.6), Inches(9), Inches(0.4),
       '不是给所有人的万能工具 — 是为情报驱动决策的团队量身定制的效率倍增器',
       sz=12, c=GRAY_400, al=PP_ALIGN.CENTER)
    pn(sl, 15)

def s16(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '部署启动', APPLE_BLUE, 16)
    steps = [
        ('Step 1', 'docker compose up -d', '一条命令启动全部服务', APPLE_BLUE),
        ('Step 2', '浏览器访问 :8765', '登录系统\n开始使用', APPLE_GREEN),
        ('Step 3', '配置数据源 + 项目', '开始自动采集\nAI Agent 立即工作', APPLE_BLUE),
    ]
    for i, (step, cmd, desc, clr) in enumerate(steps):
        l = Inches(0.5 + i*3.2)
        apple_card(sl, l, Inches(1.3), Inches(2.9), Inches(2.6), bd=clr)
        apple_bar(sl, l+Inches(0.3), Inches(1.5), Inches(0.7), Pt(3), clr)
        tx(sl, l+Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.5), step, sz=16, bold=True, c=clr, al=PP_ALIGN.CENTER)
        tx(sl, l+Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.5), cmd, sz=11, c=APPLE_BLUE, al=PP_ALIGN.CENTER, bold=True)
        tx(sl, l+Inches(0.2), Inches(3.1), Inches(2.5), Inches(0.7), desc, sz=10, c=GRAY_600, al=PP_ALIGN.CENTER)
        if i < 2:
            tx(sl, l+Inches(2.85), Inches(2.4), Inches(0.4), Inches(0.4), '→', sz=18, c=GRAY_400, al=PP_ALIGN.CENTER)
    # 开箱即用
    tx(sl, Inches(0.6), Inches(4.4), Inches(4), Inches(0.4),
       '开箱即用的能力', sz=16, bold=True, c=APPLE_BLUE)
    feats = [
        '✓ 15 个功能页面，即开即用',
        '✓ RBAC 权限体系，三级角色即刻生效',
        '✓ AI Agent 预设模板，配置即用',
        '✓ 暗色模式、响应式设计，现代用户体验',
        '✓ 新业务域：一份配置 + 一个模板 → 数天内上线',
    ]
    for i, f in enumerate(feats):
        tx(sl, Inches(0.6), Inches(4.9 + i*0.42), Inches(9), Inches(0.35),
           f, sz=11, c=GRAY_800)
    pn(sl, 16)

def s17(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG); hdr(sl, '投资回报', APPLE_GREEN, 17)
    # 三个大数字
    nums = [
        ('1-3', '个月\n投资回报周期', APPLE_GREEN),
        ('50-100', '万\n年新增收入', APPLE_BLUE),
        ('0.5-1', '个\n全职人力释放', APPLE_BLUE),
    ]
    for i, (num, sfx, clr) in enumerate(nums):
        x = Inches(0.5 + i*3.2)
        apple_card(sl, x, Inches(1.3), Inches(2.9), Inches(2.4), bd=clr)
        apple_bar(sl, x+Inches(0.3), Inches(1.5), Inches(0.7), Pt(3), clr)
        bn(sl, x+Inches(0.1), Inches(1.8), Inches(2.7), Inches(1.2), num, sfx=sfx, c=clr, fs=40)
    # 基准
    apple_card(sl, Inches(0.6), Inches(4.1), Inches(8.8), Inches(0.5))
    tx(sl, Inches(0.9), Inches(4.2), Inches(8.2), Inches(0.35),
       '基准假设：一名年薪资 20 万的销售人员', sz=11, c=GRAY_600)
    # 详细数据
    data = [
        ('人力释放', 'AI 替代 1-2 小时/天 → 相当于 0.5-1 个全职'),
        ('新增商机收入', '转化率提升 20% + 销售周期缩短 20% → 额外 50-100 万年收入'),
        ('系统部署成本', '几乎为零（开源 + 自有服务器）'),
        ('投资回报周期', '1-3 个月'),
    ]
    for i, (k, v) in enumerate(data):
        y = Inches(4.9 + i*0.45)
        apple_card(sl, Inches(0.6), y, Inches(8.8), Inches(0.38), bd=GRAY_200)
        tx(sl, Inches(0.8), y+Inches(0.06), Inches(2.0), Inches(0.25), k, sz=10, bold=True, c=APPLE_BLUE)
        tx(sl, Inches(2.9), y+Inches(0.06), Inches(6.3), Inches(0.25), v, sz=9, c=GRAY_800)
    pn(sl, 17)

def s18(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, BG)
    apple_bar(sl, Inches(0), Inches(0), Inches(10), Pt(3), APPLE_BLUE)
    tx(sl, Inches(1), Inches(2.5), Inches(8), Inches(1.0),
       'Intelligence Web', sz=48, bold=True, c=GRAY_900, al=PP_ALIGN.CENTER)
    tx(sl, Inches(1), Inches(3.7), Inches(8), Inches(0.6),
       '让情报成为您的核心竞争力', sz=22, c=APPLE_BLUE, al=PP_ALIGN.CENTER)
    apple_bar(sl, Inches(4.3), Inches(4.5), Inches(1.4), Pt(2), APPLE_GREEN)
    tx(sl, Inches(2), Inches(4.9), Inches(6), Inches(0.8),
       '谢谢\nThank You',
       sz=24, c=GRAY_400, al=PP_ALIGN.CENTER)
    pn(sl, 18)

# ============================================================
# 主函数
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s1(prs); s2(prs); s3(prs); s4(prs); s5(prs); s6(prs); s7(prs); s8(prs)
    s9(prs); s10(prs); s11(prs); s12(prs); s13(prs); s14(prs); s15(prs)
    s16(prs); s17(prs); s18(prs)

    prs.save('docs/Intelligence_Web_苹果风格 PPT.pptx')
    print('PPT 苹果风格已生成：docs/Intelligence_Web_苹果风格 PPT.pptx')
    print('18 页 · 16:9 宽屏 · 苹果极简风格 · 大量留白 · 精致布局')

if __name__ == '__main__':
    main()